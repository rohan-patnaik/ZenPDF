"""Conversion utilities for the worker process."""

from __future__ import annotations

import ipaddress
import json
import logging
import math
import os
import re
import shlex
import shutil
import socket
import subprocess
import tempfile
import time
import uuid
import zipfile
from concurrent.futures import ThreadPoolExecutor, as_completed
from html.parser import HTMLParser
from io import BytesIO
from pathlib import Path
from typing import Any, Callable, Iterable, List, Sequence, Tuple
from urllib.parse import urlparse

import img2pdf
import fitz
import requests
from docx import Document
from fpdf import FPDF
from openpyxl import Workbook
from openpyxl.utils import get_column_letter
from PIL import Image, ImageChops, ImageFilter, ImageOps
from pypdf import PdfReader, PdfWriter, Transformation
from pypdf.errors import PdfReadError
from requests_toolbelt.adapters.host_header_ssl import HostHeaderSSLAdapter

try:
    import pytesseract
except ImportError:  # pragma: no cover - optional dependency for OCR tools
    pytesseract = None

try:
    from pptx import Presentation
except ImportError:  # pragma: no cover - optional dependency for PPTX output
    Presentation = None

try:
    from pptx.util import Emu, Pt
except ImportError:  # pragma: no cover - optional dependency for PPTX output
    Emu = None
    Pt = None

try:
    import pdfplumber
except ImportError:  # pragma: no cover - optional dependency for structured table extraction
    pdfplumber = None

try:
    from playwright.sync_api import Error as PlaywrightError, sync_playwright
except ImportError:  # pragma: no cover - optional dependency for hardened browser rendering
    PlaywrightError = RuntimeError
    sync_playwright = None

try:
    from pyhanko.pdf_utils.incremental_writer import IncrementalPdfFileWriter
    from pyhanko.sign import signers
except ImportError:  # pragma: no cover - optional dependency for cryptographic signatures
    IncrementalPdfFileWriter = None
    signers = None


OCR_DPI = 300
DEFAULT_OCR_LANG = os.getenv("ZENPDF_OCR_LANG", "eng")
LOGGER = logging.getLogger(__name__)

OCR_TEXT_DENSITY_THRESHOLD = float(os.getenv("ZENPDF_OCR_TEXT_DENSITY_THRESHOLD", "60"))
OCR_PREPROCESS_THRESHOLD = int(os.getenv("ZENPDF_OCR_PREPROCESS_THRESHOLD", "180"))
DEFAULT_PAGE_RANGE_MODE = os.getenv("ZENPDF_PAGE_RANGE_MODE", "strict").strip().lower()
DEFAULT_OCR_PROFILE = os.getenv("ZENPDF_OCR_PROFILE", "balanced").strip().lower()
OCR_PROFILE_CONFIG: dict[str, dict[str, int]] = {
    "fast": {"dpi": 220, "median_filter": 1, "threshold": 168},
    "balanced": {"dpi": 300, "median_filter": 3, "threshold": OCR_PREPROCESS_THRESHOLD},
    "accurate": {"dpi": 420, "median_filter": 5, "threshold": 194},
}


def _parse_ranges(
    value: str,
    total_pages: int,
    tolerant: bool = False,
) -> List[Tuple[int, int]]:
    """
    Parse a comma-separated list of page ranges strictly.

    Examples:
    - "1,3-5"
    - "2-4,6"
    """
    ranges: List[Tuple[int, int]] = []
    tokens = [part.strip() for part in value.split(",")]
    if not tokens or all(not token for token in tokens):
        return ranges

    invalid_tokens: list[str] = []
    for token in tokens:
        try:
            if not token:
                raise ValueError("Page ranges cannot contain empty items")
            if "-" in token:
                pieces = token.split("-")
                if len(pieces) != 2:
                    raise ValueError(f"Invalid page range token: {token!r}")
                start_text, end_text = pieces[0].strip(), pieces[1].strip()
            else:
                start_text = token
                end_text = token
            if not start_text or not end_text:
                raise ValueError(f"Invalid page range token: {token!r}")
            try:
                start_i = int(start_text)
                end_i = int(end_text)
            except ValueError as error:
                raise ValueError(f"Invalid page number in token: {token!r}") from error
            if start_i < 1 or end_i < 1:
                raise ValueError("Page numbers must be >= 1")
            if start_i > end_i:
                raise ValueError(f"Range start is greater than end: {token!r}")
            if start_i > total_pages:
                raise ValueError(
                    f"Page {start_i} exceeds document page count ({total_pages})"
                )
            clamped_end = min(end_i, total_pages)
            ranges.append((start_i, clamped_end))
        except ValueError:
            if not tolerant:
                raise
            invalid_tokens.append(token)
            LOGGER.warning("Ignoring invalid range token in tolerant mode: %r", token)
            continue
    if tolerant and invalid_tokens and not ranges:
        joined = ", ".join(repr(token) for token in invalid_tokens)
        raise ValueError(f"No valid page ranges found; invalid tokens: {joined}")
    return ranges


def _dedupe_preserve_order(values: Iterable[int]) -> List[int]:
    """Remove duplicate integers while preserving first-seen order."""
    deduped: List[int] = []
    seen: set[int] = set()
    for value in values:
        if value in seen:
            continue
        seen.add(value)
        deduped.append(value)
    return deduped


def _parse_page_list(
    value: str,
    total_pages: int,
    dedupe: bool = True,
    tolerant: bool = False,
) -> List[int]:
    """
    Expand a comma-separated page range string into an ordered list of page numbers.
    
    Parameters:
        value (str): A comma-separated string of page numbers and ranges (e.g. "1,3-5,7").
        total_pages (int): Total number of pages in the document; results are clamped to 1..total_pages.
    
    Returns:
        list[int]: Ordered list of page numbers produced by expanding the ranges in `value`.
                   Each page is within 1 and `total_pages` inclusive; overlapping ranges may produce duplicates.
    """
    pages: List[int] = []
    for start, end in _parse_ranges(value, total_pages, tolerant=tolerant):
        pages.extend(range(start, end + 1))
    clamped = [page for page in pages if 1 <= page <= total_pages]
    return _dedupe_preserve_order(clamped) if dedupe else clamped


def _load_pdf(input_path: Path, allow_encrypted: bool = False) -> PdfReader:
    """Load a PDF and optionally enforce unencrypted input."""
    try:
        reader = PdfReader(str(input_path))
    except PdfReadError as error:
        raise ValueError("PDF appears to be corrupted or unreadable.") from error
    if reader.is_encrypted and not allow_encrypted:
        raise ValueError("PDF is encrypted")
    return reader


def _resolve_page_selection(
    pages: str | None,
    total_pages: int,
    tolerant: bool | None = None,
) -> set[int] | None:
    """Return a validated set of target pages or None for all."""
    if pages is None:
        return None
    value = str(pages).strip()
    if not value:
        return None
    tolerant_mode = (
        DEFAULT_PAGE_RANGE_MODE == "tolerant" if tolerant is None else bool(tolerant)
    )
    selection = _parse_page_list(value, total_pages, tolerant=tolerant_mode)
    if not selection:
        raise ValueError("No valid pages selected")
    return set(selection)


def _parse_bool(value: Any, default: bool = False) -> bool:
    """Parse booleans from string/int forms."""
    if value is None:
        return default
    if isinstance(value, bool):
        return value
    if isinstance(value, (int, float)):
        return bool(value)
    normalized = str(value).strip().lower()
    if normalized in {"1", "true", "yes", "on", "y"}:
        return True
    if normalized in {"0", "false", "no", "off", "n"}:
        return False
    return default


def _normalize_ocr_profile(profile: str | None) -> str:
    """Return a supported OCR profile name."""
    normalized = (profile or DEFAULT_OCR_PROFILE).strip().lower()
    if normalized not in OCR_PROFILE_CONFIG:
        return "balanced"
    return normalized


def _ocr_profile_settings(profile: str | None) -> dict[str, int]:
    """Return OCR preprocessing settings for a profile."""
    return OCR_PROFILE_CONFIG[_normalize_ocr_profile(profile)]


def _safe_text_length(page: Any) -> int:
    """Return extracted text length for a page without raising."""
    try:
        return len((page.extract_text() or "").strip())
    except Exception:
        return 0


def _pdf_text_density(input_path: Path) -> float:
    """Estimate average extracted text length per page."""
    reader = _load_pdf(input_path)
    if not reader.pages:
        return 0.0
    total_chars = sum(_safe_text_length(page) for page in reader.pages)
    return total_chars / max(1, len(reader.pages))


def _is_text_light_pdf(input_path: Path) -> bool:
    """Heuristic to decide if OCR assistance should be preferred."""
    return _pdf_text_density(input_path) < OCR_TEXT_DENSITY_THRESHOLD


def _copy_metadata(writer: PdfWriter, reader: PdfReader) -> None:
    """Copy metadata from a PDF reader into a writer."""
    metadata = reader.metadata or {}
    if metadata:
        writer.add_metadata(metadata)


def _assert_fitz_unencrypted(document: fitz.Document) -> None:
    """Raise if a PyMuPDF document is encrypted."""
    is_encrypted = bool(
        getattr(document, "is_encrypted", False)
        or getattr(document, "isEncrypted", False)
    )
    if is_encrypted:
        raise ValueError("PDF is encrypted")


def _parse_margins(value: str) -> Tuple[float, float, float, float]:
    """
    Parse a comma-separated margin string into top, right, bottom, and left values in points.
    
    The input may be a single numeric value (applied to all four margins) or four comma-separated numeric values in the order top, right, bottom, left. Whitespace around values is ignored.
    
    Parameters:
        value (str): Margin specification as "N" or "T,R,B,L".
    
    Returns:
        tuple[float, float, float, float]: (top, right, bottom, left) in points.
    
    Raises:
        ValueError: If any component is not numeric or if the input does not contain 1 or 4 values.
    """
    parts = [part.strip() for part in value.split(",") if part.strip()]
    try:
        numbers = [float(part) for part in parts]
    except ValueError as error:
        raise ValueError("Margins must be numeric") from error
    if len(numbers) == 1:
        top = right = bottom = left = numbers[0]
    elif len(numbers) == 4:
        top, right, bottom, left = numbers
    else:
        raise ValueError("Margins must have 1 or 4 values")
    return (top, right, bottom, left)


def _rotate_page(page, angle: int) -> None:
    """
    Rotate the given PDF page object by the specified angle in degrees, modifying the page in place.
    
    Parameters:
        page: PDF page object that provides a rotation method (one of `rotate_clockwise`, `rotateClockwise`, or `rotate`).
        angle (int): Rotation angle in degrees clockwise.
    """
    if hasattr(page, "rotate_clockwise"):
        page.rotate_clockwise(angle)
    elif hasattr(page, "rotateClockwise"):
        page.rotateClockwise(angle)
    elif hasattr(page, "rotate"):
        page.rotate(angle)


def _points_to_mm(points: float) -> float:
    """
    Convert a length in PDF points to millimeters.
    
    Returns:
        millimeters (float): The length converted from points to millimeters.
    """
    return points * 25.4 / 72


def _build_overlay_page(
    width_points: float,
    height_points: float,
    draw_fn: Callable[[FPDF, float, float], None],
):
    """
    Create a single-page PDF overlay sized to the given page dimensions and rendered by the provided draw callback.
    
    Parameters:
        width_points (float): Page width in PDF points.
        height_points (float): Page height in PDF points.
        draw_fn (Callable[[FPDF, float, float], None]): Callback that draws onto an FPDF instance; called with the FPDF object and the page width and height in millimeters.
    
    Returns:
        page: A single page object from a PdfReader representing the generated overlay.
    """
    width_mm = _points_to_mm(width_points)
    height_mm = _points_to_mm(height_points)
    orientation = "L" if width_mm > height_mm else "P"
    pdf = FPDF(orientation=orientation, unit="mm", format=(width_mm, height_mm))
    pdf.set_margins(0, 0, 0)
    pdf.set_auto_page_break(auto=False)
    pdf.add_page()
    draw_fn(pdf, width_mm, height_mm)
    pdf_output = pdf.output()
    if isinstance(pdf_output, str):
        pdf_bytes = pdf_output.encode("latin-1")
    else:
        pdf_bytes = bytes(pdf_output)
    overlay_reader = PdfReader(BytesIO(pdf_bytes))
    return overlay_reader.pages[0]


def _merge_overlay_page(
    page: Any,
    overlay: Any,
    box: Any,
) -> None:
    """
    Merge an overlay page onto a target PDF page, honoring crop box offsets when present.
    
    If the crop box origin is not at (0, 0), translates the overlay so it aligns with the visible page region.
    Falls back to a direct merge if transformed merge is unavailable.
    """
    try:
        lower_left = box.lower_left
        tx = float(lower_left[0])
        ty = float(lower_left[1])
    except Exception:
        tx = 0.0
        ty = 0.0

    if (tx or ty) and hasattr(page, "merge_transformed_page"):
        page.merge_transformed_page(overlay, Transformation().translate(tx, ty))
    else:
        page.merge_page(overlay)


def merge_pdfs(inputs: Sequence[Path], output_path: Path) -> Path:
    """
    Merge multiple PDF files into a single PDF.
    
    Parameters:
        inputs (Sequence[Path]): Paths to source PDF files, merged in the given order.
        output_path (Path): Destination path for the merged PDF.
    
    Returns:
        Path: The path to the written merged PDF (same as `output_path`).
    """
    writer = PdfWriter()
    for path in inputs:
        reader = PdfReader(str(path))
        for page in reader.pages:
            writer.add_page(page)
    with output_path.open("wb") as handle:
        writer.write(handle)
    return output_path


def split_pdf(
    input_path: Path,
    output_dir: Path,
    ranges: str | None,
    tolerant_ranges: bool = False,
) -> List[Path]:
    """Split a PDF into multiple files based on ranges."""
    reader = PdfReader(str(input_path))
    total_pages = len(reader.pages)
    output_files: List[Path] = []

    if ranges:
        page_ranges = _parse_ranges(ranges, total_pages, tolerant=tolerant_ranges)
        if not page_ranges:
            raise ValueError("No valid page ranges provided")
    else:
        page_ranges = [(index, index) for index in range(1, total_pages + 1)]

    for index, (start, end) in enumerate(page_ranges, start=1):
        writer = PdfWriter()
        for page_number in range(start - 1, end):
            writer.add_page(reader.pages[page_number])
        output_path = output_dir / f"split_{index}.pdf"
        with output_path.open("wb") as handle:
            writer.write(handle)
        output_files.append(output_path)
    return output_files


def compress_pdf(input_path: Path, output_path: Path) -> tuple[Path, dict]:
    """Compress a PDF using a staged pipeline with early image-heavy handling."""
    size_bytes = input_path.stat().st_size
    size_mb = max(1, math.ceil(size_bytes / (1024 * 1024)))
    warnings: list[str] = []
    steps: list[dict] = []
    candidates: list[dict] = []
    run_id = uuid.uuid4().hex[:8]

    def _env_int(name: str, default: int) -> int:
        value = os.environ.get(name)
        if value is None or value == "":
            return default
        try:
            return int(value)
        except ValueError:
            return default

    def _env_float(name: str, default: float) -> float:
        value = os.environ.get(name)
        if value is None or value == "":
            return default
        try:
            return float(value)
        except ValueError:
            return default

    def _env_bool(name: str, default: bool = False) -> bool:
        value = os.environ.get(name)
        if value is None:
            return default
        return value.strip().lower() in {"1", "true", "yes", "on", "y"}

    def _env_str(name: str, default: str) -> str:
        value = os.environ.get(name)
        return value.strip() if value else default

    def _run_cmd(cmd: list[str], timeout_s: int, env: dict | None = None) -> dict:
        start = time.perf_counter()
        try:
            result = subprocess.run(
                cmd,
                capture_output=True,
                text=True,
                check=False,
                timeout=timeout_s,
                env=env,
            )
            elapsed_ms = int((time.perf_counter() - start) * 1000)
            return {
                "ok": result.returncode == 0,
                "returncode": result.returncode,
                "stdout": (result.stdout or "").strip(),
                "stderr": (result.stderr or "").strip(),
                "timeout": False,
                "ms": elapsed_ms,
            }
        except subprocess.TimeoutExpired:
            elapsed_ms = int((time.perf_counter() - start) * 1000)
            return {
                "ok": False,
                "returncode": None,
                "stdout": "",
                "stderr": "",
                "timeout": True,
                "ms": elapsed_ms,
            }

    def _truncate(value: str, limit: int = 300) -> str:
        if len(value) <= limit:
            return value
        return f"{value[:limit]}..."

    def _make_step_entry(name: str, result: dict | None, notes: str | None = None) -> dict:
        if result is None:
            return {"name": name, "ok": False, "ms": 0, "notes": notes or ""}
        entry = {"name": name, "ok": bool(result.get("ok")), "ms": int(result.get("ms", 0))}
        if notes:
            entry["notes"] = _truncate(notes)
        elif result.get("timeout"):
            entry["notes"] = "timeout"
        elif not result.get("ok") and result.get("stderr"):
            entry["notes"] = _truncate((result.get("stderr") or "").strip())
        return entry

    def _record_step(name: str, result: dict | None, notes: str | None = None) -> None:
        steps.append(_make_step_entry(name, result, notes))

    def _safe_page_count(path: Path) -> int:
        try:
            reader = PdfReader(str(path))
            if reader.is_encrypted:
                raise ValueError("PDF is encrypted")
            return max(1, len(reader.pages))
        except ValueError:
            raise
        except (PdfReadError, OSError, EOFError) as error:
            warnings.append(f"Could not count pages: {error}")
        except Exception:
            pass
        try:
            with fitz.open(str(path)) as document:
                return max(1, document.page_count)
        except Exception:
            return 1

    def _validate_candidate(path: Path, expected_pages: int) -> bool:
        if not path.exists() or path.stat().st_size == 0:
            return False
        if shutil.which("qpdf"):
            result = _run_cmd(["qpdf", "--check", str(path)], timeout_s=20)
            if not result["ok"]:
                return False
        try:
            with fitz.open(str(path)) as document:
                if document.page_count != expected_pages:
                    return False
                if document.page_count > 0:
                    page = document.load_page(0)
                    _ = page.get_pixmap(matrix=fitz.Matrix(0.5, 0.5))
        except Exception:
            return False
        return True

    def _add_candidate(path: Path, method: str, label: str, expected_pages: int) -> bool:
        if not _validate_candidate(path, expected_pages):
            warnings.append(f"{label} output invalid")
            return False
        candidates.append(
            {
                "path": path,
                "method": method,
                "label": label,
                "size": path.stat().st_size,
            }
        )
        return True

    def _rewrite_pdf(source: Path, target: Path) -> bool:
        try:
            reader = PdfReader(str(source))
            if reader.is_encrypted:
                raise ValueError("PDF is encrypted")
            writer = PdfWriter()
            for page in reader.pages:
                writer.add_page(page)
            compress = getattr(writer, "compress_content_streams", None)
            if callable(compress):
                compress()
            raw_metadata = reader.metadata or {}
            safe_metadata = {}
            for key, value in raw_metadata.items():
                if value is None:
                    continue
                safe_metadata[str(key)] = str(value)
            if safe_metadata:
                writer.add_metadata(safe_metadata)
            with target.open("wb") as handle:
                writer.write(handle)
            return True
        except Exception as error:
            warnings.append(f"pypdf rewrite failed: {error}")
            return False

    def _detect_image_heavy(path: Path) -> dict:
        metrics = {
            "pages": 0,
            "sampled_pages": 0,
            "image_count": 0,
            "text_chars": 0,
            "image_heavy": False,
        }
        try:
            with fitz.open(str(path)) as document:
                pages = document.page_count
                metrics["pages"] = pages
                if pages == 0:
                    return metrics
                sample_pages = min(pages, 10)
                step = max(1, pages // sample_pages)
                sampled = 0
                image_count = 0
                text_chars = 0
                for index in range(0, pages, step):
                    if sampled >= sample_pages:
                        break
                    page = document.load_page(index)
                    image_count += len(page.get_images(full=True))
                    text_chars += len(page.get_text("text") or "")
                    sampled += 1
                metrics["sampled_pages"] = sampled
                if sampled > 0:
                    images_per_page = image_count / sampled
                    text_per_page = text_chars / sampled
                else:
                    images_per_page = 0
                    text_per_page = 0
                metrics["image_count"] = int(round(images_per_page * pages))
                metrics["text_chars"] = int(round(text_per_page * pages))
                metrics["images_per_page"] = images_per_page
                metrics["text_chars_per_page"] = text_per_page
                image_heavy = metrics["image_count"] >= pages * 1.0 or (
                    (text_per_page < 500) and (images_per_page > 0.5)
                )
                metrics["image_heavy"] = image_heavy
        except Exception as error:
            warnings.append(f"image-heavy detector failed: {error}")
        return metrics

    def _tmp_path(name: str) -> Path:
        return output_path.parent / f"{input_path.stem}_{run_id}_{name}"

    try:
        _load_pdf(input_path, allow_encrypted=False)
    except ValueError as error:
        if str(error) == "PDF is encrypted":
            raise
        warnings.append(f"preflight read failed: {error}")
    except Exception as error:  # noqa: BLE001
        warnings.append(f"preflight read failed: {error}")

    timeout_override = _env_int("ZENPDF_COMPRESS_TIMEOUT_SECONDS", 0)
    if timeout_override > 0:
        timeout_seconds = timeout_override
    else:
        base_timeout = _env_int("ZENPDF_COMPRESS_TIMEOUT_BASE_SECONDS", 120)
        per_mb_timeout = _env_int("ZENPDF_COMPRESS_TIMEOUT_PER_MB_SECONDS", 3)
        per_page_timeout = _env_float("ZENPDF_COMPRESS_TIMEOUT_PER_PAGE_SECONDS", 1.5)
        max_timeout = _env_int("ZENPDF_COMPRESS_TIMEOUT_MAX_SECONDS", 900)
        pages_for_timeout = _safe_page_count(input_path)
        timeout_seconds = min(
            max_timeout,
            int(base_timeout + (size_mb * per_mb_timeout) + (pages_for_timeout * per_page_timeout)),
        )

    probe_pages = max(1, min(_safe_page_count(input_path), _env_int("ZENPDF_COMPRESS_TIMEOUT_PROBE_PAGES", 5)))
    probe_timeout = min(
        _env_int("ZENPDF_COMPRESS_TIMEOUT_PROBE_MAX_SECONDS", 30),
        max(10, int(timeout_seconds * 0.25)),
    )

    profile = _env_str("ZENPDF_COMPRESS_PROFILE", "balanced").lower()
    if profile not in {"light", "balanced", "strong"}:
        profile = "balanced"

    auto_image_heavy = _env_bool("ZENPDF_COMPRESS_AUTO_IMAGE_HEAVY", True)
    pass_through_env = _env_bool("ZENPDF_COMPRESS_GS_PASSTHROUGH_JPEG", False)
    use_zopfli = _env_bool("ZENPDF_COMPRESS_USE_ZOPFLI", False)
    enable_image_opt = _env_bool("ZENPDF_COMPRESS_ENABLE_IMAGE_OPT", False)
    enable_pdfsizeopt = _env_bool("ZENPDF_COMPRESS_ENABLE_PDFSIZEOPT", False)
    enable_jbig2 = _env_bool("ZENPDF_COMPRESS_ENABLE_JBIG2", False)
    qpdf_keep_inline = _env_bool("ZENPDF_QPDF_OI_KEEP_INLINE_IMAGES", False)
    gs_min_size_mb = _env_int("ZENPDF_COMPRESS_GS_MIN_SIZE_MB", 5)
    gs_preset = _env_str("ZENPDF_COMPRESS_GS_PRESET", "").lower()
    gs_extra_flags = _env_bool("ZENPDF_COMPRESS_GS_EXTRA_FLAGS", False)
    mutool_object_streams = _env_bool("ZENPDF_MUTOOL_OBJECT_STREAMS", False)
    parallelism = max(1, _env_int("ZENPDF_COMPRESS_PARALLELISM", 1))
    qpdf_quality = _env_int("ZENPDF_QPDF_OI_QUALITY", 75)
    qpdf_min_width = _env_int("ZENPDF_QPDF_OI_MIN_WIDTH", 128)
    qpdf_min_height = _env_int("ZENPDF_QPDF_OI_MIN_HEIGHT", 128)
    qpdf_min_area = _env_int("ZENPDF_QPDF_OI_MIN_AREA", 16384)
    pdfsizeopt_args = shlex.split(_env_str("ZENPDF_COMPRESS_PDFSIZEOPT_ARGS", ""))

    savings_threshold_env = os.environ.get("ZENPDF_COMPRESS_SAVINGS_THRESHOLD_PCT")
    if savings_threshold_env:
        try:
            savings_threshold = float(savings_threshold_env)
        except ValueError:
            savings_threshold = 0.08
    else:
        legacy_percent = os.environ.get("ZENPDF_COMPRESS_MIN_SAVINGS_PERCENT")
        if legacy_percent:
            try:
                savings_threshold = float(legacy_percent) / 100.0
            except ValueError:
                savings_threshold = 0.08
        else:
            savings_threshold = 0.08
    min_savings_bytes = _env_int("ZENPDF_COMPRESS_MIN_SAVINGS_BYTES", 200000)

    expected_pages = _safe_page_count(input_path)
    image_metrics = _detect_image_heavy(input_path)
    image_heavy = auto_image_heavy and image_metrics.get("image_heavy", False)

    mutool = shutil.which("mutool")
    qpdf = shutil.which("qpdf")
    ghostscript = shutil.which("gs")
    pdfsizeopt = shutil.which("pdfsizeopt")
    jbig2 = shutil.which("jbig2")

    if _validate_candidate(input_path, expected_pages):
        _add_candidate(input_path, "original", "original", expected_pages)

    normalized_path = input_path
    base_path = input_path
    gs_input_path = input_path
    temp_paths: list[Path] = []

    if mutool:
        normalized_path = _tmp_path("normalized.pdf")
        cmd = ["mutool", "clean", "-gggg", "-z", "-i", "-f", "-t"]
        if mutool_object_streams:
            cmd.append("-Z")
        cmd.extend([str(input_path), str(normalized_path)])
        result = _run_cmd(cmd, timeout_seconds)
        _record_step("normalize_mutool", result)
        temp_paths.append(normalized_path)
        if result["ok"]:
            _add_candidate(normalized_path, "mutool", "normalize", expected_pages)
            base_path = normalized_path
            gs_input_path = normalized_path
    else:
        _record_step("normalize_mutool", None, "skipped: mutool not available")

    if normalized_path == input_path:
        if qpdf:
            qpdf_norm = _tmp_path("qpdf_norm.pdf")
            cmd = [
                "qpdf",
                "--warning-exit-0",
                "--object-streams=generate",
                "--compress-streams=y",
                "--recompress-flate",
                str(input_path),
                str(qpdf_norm),
            ]
            result = _run_cmd(cmd, timeout_seconds)
            _record_step("normalize_qpdf", result)
            temp_paths.append(qpdf_norm)
            if result["ok"]:
                normalized_path = qpdf_norm
                _add_candidate(qpdf_norm, "qpdf", "normalize", expected_pages)
                base_path = qpdf_norm
                gs_input_path = qpdf_norm
        else:
            _record_step("normalize_qpdf", None, "skipped: qpdf not available")

    if normalized_path == input_path:
        rewrite_path = _tmp_path("rewrite.pdf")
        start = time.perf_counter()
        ok = _rewrite_pdf(input_path, rewrite_path)
        elapsed_ms = int((time.perf_counter() - start) * 1000)
        result = {"ok": ok, "ms": elapsed_ms, "stderr": "" if ok else "rewrite failed"}
        _record_step("normalize_pypdf", result)
        temp_paths.append(rewrite_path)
        if ok:
            normalized_path = rewrite_path
            _add_candidate(rewrite_path, "pypdf", "normalize", expected_pages)
            base_path = rewrite_path
            gs_input_path = rewrite_path

    if qpdf:
        optimized_path = _tmp_path("optimized.pdf")
        cmd = [
            "qpdf",
            "--warning-exit-0",
            "--object-streams=generate",
            "--compress-streams=y",
            "--recompress-flate",
            str(base_path),
            str(optimized_path),
        ]
        result = _run_cmd(cmd, timeout_seconds)
        _record_step("optimize_qpdf", result)
        temp_paths.append(optimized_path)
        if result["ok"]:
            _add_candidate(optimized_path, "qpdf", "optimize", expected_pages)
            base_path = optimized_path
            gs_input_path = optimized_path
    else:
        _record_step("optimize_qpdf", None, "skipped: qpdf not available")

    if mutool:
        mutool_opt_path = _tmp_path("mutool_opt.pdf")
        cmd = [
            "mutool",
            "merge",
            "-o",
            str(mutool_opt_path),
            "-O",
            "compress",
            str(base_path),
        ]
        result = _run_cmd(cmd, timeout_seconds)
        _record_step("optimize_mutool", result)
        temp_paths.append(mutool_opt_path)
        if result["ok"]:
            _add_candidate(mutool_opt_path, "mutool", "mutool_opt", expected_pages)
    else:
        _record_step("optimize_mutool", None, "skipped: mutool not available")

    def _ghostscript_settings(override_preset: str | None = None) -> tuple[str, str, str]:
        preset = override_preset or ("screen" if profile == "strong" else "ebook")
        if preset == "screen":
            return "/screen", "100", "100"
        return "/ebook", "150", "150"

    def _run_gs_basic(source: Path, dest: Path, label: str, preset: str | None = None) -> dict:
        if not ghostscript:
            return {"ok": False, "stderr": "ghostscript missing", "timeout": False, "ms": 0}
        pdfsettings, color_res, gray_res = _ghostscript_settings(preset)
        pass_through = pass_through_env or profile == "light"
        cmd = [
            ghostscript,
            "-dSAFER",
            "-dBATCH",
            "-dNOPAUSE",
            "-sDEVICE=pdfwrite",
            "-dCompatibilityLevel=1.5",
            f"-dPDFSETTINGS={pdfsettings}",
            f"-dPassThroughJPEGImages={'true' if pass_through else 'false'}",
            "-dDownsampleColorImages=true",
            f"-dColorImageResolution={color_res}",
            "-dColorImageDownsampleType=/Bicubic",
            "-dDownsampleGrayImages=true",
            f"-dGrayImageResolution={gray_res}",
            "-dGrayImageDownsampleType=/Bicubic",
            "-dDownsampleMonoImages=true",
            "-dMonoImageResolution=300",
            "-dAutoFilterColorImages=false",
            "-dColorImageFilter=/DCTEncode",
            "-dAutoFilterGrayImages=false",
            "-dGrayImageFilter=/DCTEncode",
        ]
        if gs_extra_flags:
            cmd.append("-dDetectDuplicateImages=true")
        cmd.extend([f"-sOutputFile={dest}", str(source)])
        result = _run_cmd(cmd, timeout_seconds)
        _record_step(label, result)
        return result

    if image_heavy and ghostscript and size_mb >= gs_min_size_mb:
        gs_output = _tmp_path("gs_early.pdf")
        temp_paths.append(gs_output)
        preset_override = "screen" if profile == "strong" else "ebook"
        result = _run_gs_basic(normalized_path, gs_output, "ghostscript_early", preset_override)
        if result["ok"]:
            _add_candidate(gs_output, "ghostscript", "ghostscript_early", expected_pages)
            if qpdf:
                qpdf_after_gs = _tmp_path("gs_qpdf.pdf")
                temp_paths.append(qpdf_after_gs)
                result = _run_cmd(
                    [
                        "qpdf",
                        "--object-streams=generate",
                        "--compress-streams=y",
                        "--recompress-flate",
                        "--compression-level=9",
                        "--deterministic-id",
                        str(gs_output),
                        str(qpdf_after_gs),
                    ],
                    timeout_seconds,
                )
                _record_step("qpdf_after_gs", result)
                if result["ok"]:
                    _add_candidate(qpdf_after_gs, "qpdf_after_gs", "qpdf_after_gs", expected_pages)

    def _current_best_savings() -> tuple[int, float]:
        if not candidates or size_bytes == 0:
            return 0, 0.0
        best = min(candidates, key=lambda item: item["size"])
        savings = max(size_bytes - best["size"], 0)
        return savings, savings / size_bytes

    def _should_run_heavy_steps() -> bool:
        savings_bytes, savings_pct = _current_best_savings()
        return savings_bytes < min_savings_bytes or savings_pct < savings_threshold

    def _empty_outcome() -> dict:
        return {"steps": [], "candidates": [], "temp_paths": [], "warnings": []}

    def _run_image_opt_task() -> dict:
        outcome = _empty_outcome()
        if not enable_image_opt:
            return outcome
        if not qpdf:
            outcome["steps"].append(_make_step_entry("optimize_images_qpdf", None, "skipped: qpdf not available"))
            return outcome
        image_opt_path = _tmp_path("image_opt.pdf")
        outcome["temp_paths"].append(image_opt_path)
        cmd = [
            "qpdf",
            "--warning-exit-0",
            "--optimize-images",
            f"--oi-quality={qpdf_quality}",
            f"--oi-min-width={qpdf_min_width}",
            f"--oi-min-height={qpdf_min_height}",
            f"--oi-min-area={qpdf_min_area}",
            "--object-streams=generate",
            "--compress-streams=y",
            "--recompress-flate",
            "--compression-level=9",
            "--deterministic-id",
            str(base_path),
            str(image_opt_path),
        ]
        if qpdf_keep_inline:
            cmd.insert(3, "--keep-inline-images")
        result = _run_cmd(cmd, timeout_seconds)
        outcome["steps"].append(_make_step_entry("optimize_images_qpdf", result))
        if result["ok"]:
            outcome["candidates"].append((image_opt_path, "qpdf_optimize_images", "image_opt"))
        return outcome

    def _run_pdfsizeopt_task(should_run: bool) -> dict:
        outcome = _empty_outcome()
        if not (enable_pdfsizeopt or enable_jbig2):
            return outcome
        if enable_jbig2 and not jbig2:
            outcome["steps"].append(_make_step_entry("optimize_pdfsizeopt", None, "skipped: jbig2enc not available"))
            return outcome
        if not pdfsizeopt:
            outcome["steps"].append(_make_step_entry("optimize_pdfsizeopt", None, "skipped: pdfsizeopt not available"))
            return outcome
        if not should_run:
            outcome["steps"].append(_make_step_entry("optimize_pdfsizeopt", None, "skipped: already reduced"))
            return outcome
        pdfsizeopt_path = _tmp_path("pdfsizeopt.pdf")
        outcome["temp_paths"].append(pdfsizeopt_path)
        cmd = [pdfsizeopt]
        if enable_jbig2:
            cmd.append("--use-image-optimizer=jbig2")
        if pdfsizeopt_args:
            cmd.extend(pdfsizeopt_args)
        cmd.extend([str(base_path), str(pdfsizeopt_path)])
        result = _run_cmd(cmd, timeout_seconds)
        outcome["steps"].append(_make_step_entry("optimize_pdfsizeopt", result))
        if result["ok"]:
            method = "pdfsizeopt_jbig2" if enable_jbig2 else "pdfsizeopt"
            outcome["candidates"].append((pdfsizeopt_path, method, "pdfsizeopt"))
        return outcome

    def _run_ghostscript_task(should_run: bool) -> dict:
        outcome = _empty_outcome()
        if image_heavy:
            outcome["steps"].append(_make_step_entry("ghostscript_full", None, "skipped: image-heavy handled"))
            return outcome
        if not ghostscript:
            outcome["steps"].append(_make_step_entry("ghostscript_full", None, "skipped: ghostscript not available"))
            return outcome
        if size_mb < gs_min_size_mb:
            outcome["steps"].append(_make_step_entry("ghostscript_full", None, "skipped: below size threshold"))
            return outcome
        if not should_run:
            outcome["steps"].append(_make_step_entry("ghostscript_full", None, "skipped: already reduced"))
            return outcome

        preset_override = None
        if gs_preset in {"screen", "ebook"}:
            preset_override = gs_preset

        pdfsettings, color_res, gray_res = _ghostscript_settings(preset_override)
        pass_through = pass_through_env or profile == "light"
        gs_flags = [
            "-dSAFER",
            "-dBATCH",
            "-dNOPAUSE",
            "-sDEVICE=pdfwrite",
            "-dCompatibilityLevel=1.5",
            f"-dPDFSETTINGS={pdfsettings}",
            f"-dPassThroughJPEGImages={'true' if pass_through else 'false'}",
            "-dDownsampleColorImages=true",
            f"-dColorImageResolution={color_res}",
            "-dColorImageDownsampleType=/Bicubic",
            "-dDownsampleGrayImages=true",
            f"-dGrayImageResolution={gray_res}",
            "-dGrayImageDownsampleType=/Bicubic",
            "-dDownsampleMonoImages=true",
            "-dMonoImageResolution=300",
            "-dAutoFilterColorImages=false",
            "-dColorImageFilter=/DCTEncode",
            "-dAutoFilterGrayImages=false",
            "-dGrayImageFilter=/DCTEncode",
        ]
        if gs_extra_flags:
            gs_flags.append("-dDetectDuplicateImages=true")

        def _gs_cmd(
            source: Path,
            target: Path,
            first_page: int | None = None,
            last_page: int | None = None,
            newpdf: bool | None = None,
        ) -> list[str]:
            cmd = [ghostscript, *gs_flags]
            if first_page is not None and last_page is not None:
                cmd.extend([f"-dFirstPage={first_page}", f"-dLastPage={last_page}"])
            if newpdf is False:
                cmd.append("-dNEWPDF=false")
            cmd.append(f"-sOutputFile={target}")
            cmd.append(str(source))
            return cmd

        probe_output = _tmp_path("gs_probe.pdf")
        gs_output = _tmp_path("gs.pdf")
        outcome["temp_paths"].extend([probe_output, gs_output])

        probe_result = _run_cmd(
            _gs_cmd(gs_input_path, probe_output, 1, probe_pages),
            probe_timeout,
        )
        probe_ok = probe_result["ok"]
        probe_notes = None
        if probe_ok and probe_pages > 0:
            estimated_ms = int((probe_result["ms"] / probe_pages) * expected_pages)
            if estimated_ms > timeout_seconds * 1000:
                probe_ok = False
                probe_notes = "probe too slow, skipping full run"
        outcome["steps"].append(_make_step_entry("ghostscript_probe", probe_result, probe_notes))
        if probe_ok:
            full_result = _run_cmd(_gs_cmd(gs_input_path, gs_output), timeout_seconds)
            outcome["steps"].append(_make_step_entry("ghostscript_full", full_result))
            if full_result["ok"]:
                outcome["candidates"].append((gs_output, "ghostscript", "ghostscript"))
            else:
                retry_result = _run_cmd(_gs_cmd(gs_input_path, gs_output, newpdf=False), timeout_seconds)
                outcome["steps"].append(_make_step_entry("ghostscript_retry", retry_result))
                if retry_result["ok"]:
                    outcome["candidates"].append((gs_output, "ghostscript", "ghostscript_retry"))
        else:
            outcome["steps"].append(_make_step_entry("ghostscript_full", None, "skipped: probe failed"))
        return outcome

    should_run_heavy = _should_run_heavy_steps()

    task_specs: list[tuple[str, Callable[[], dict]]] = []
    if enable_image_opt:
        task_specs.append(("image_opt", _run_image_opt_task))
    if enable_pdfsizeopt or enable_jbig2:
        task_specs.append(("pdfsizeopt", lambda: _run_pdfsizeopt_task(should_run_heavy)))
    if ghostscript:
        task_specs.append(("ghostscript", lambda: _run_ghostscript_task(should_run_heavy)))

    task_results: dict[str, dict] = {}
    if parallelism > 1 and len(task_specs) > 1:
        max_workers = min(parallelism, len(task_specs))
        with ThreadPoolExecutor(max_workers=max_workers) as executor:
            future_map = {executor.submit(task): name for name, task in task_specs}
            for future in as_completed(future_map):
                name = future_map[future]
                try:
                    task_results[name] = future.result()
                except Exception as error:  # noqa: BLE001
                    task_results[name] = {
                        "steps": [
                            _make_step_entry(
                                name, {"ok": False, "ms": 0, "stderr": str(error)}
                            )
                        ],
                        "candidates": [],
                        "temp_paths": [],
                        "warnings": [f"{name} failed: {error}"],
                    }
    else:
        for name, task in task_specs:
            task_results[name] = task()

    for name, _ in task_specs:
        outcome = task_results.get(name, _empty_outcome())
        steps.extend(outcome.get("steps", []))
        for warning in outcome.get("warnings", []):
            warnings.append(warning)
        for path, method, label in outcome.get("candidates", []):
            _add_candidate(path, method, label, expected_pages)
        temp_paths.extend(outcome.get("temp_paths", []))

    if not candidates:
        raise ValueError(
            "Could not compress this PDF due to malformed structure; try Repair PDF first."
        )

    best = min(candidates, key=lambda item: item["size"])
    method = best["method"]
    if method == "original":
        method = "passthrough"
        warnings.append("No smaller output found; preserving original content.")

    savings_bytes = max(size_bytes - best["size"], 0)
    savings_pct = (savings_bytes / size_bytes) if size_bytes else 0.0

    status = "success"
    if savings_bytes < min_savings_bytes or savings_pct < savings_threshold:
        status = "no_change"
        best = next((c for c in candidates if c["path"] == input_path), best)

    final_path = best["path"]
    if final_path != output_path:
        if final_path == input_path:
            shutil.copy2(final_path, output_path)
        else:
            if output_path.exists():
                output_path.unlink()
            final_path.replace(output_path)
        final_path = output_path

    if qpdf:
        deterministic_out = _tmp_path("deterministic.pdf")
        result = _run_cmd(
            [
                "qpdf",
                "--object-streams=generate",
                "--compress-streams=y",
                "--recompress-flate",
                "--compression-level=9",
                "--deterministic-id",
                str(final_path),
                str(deterministic_out),
            ],
            min(timeout_seconds, 120),
        )
        _record_step("deterministic_id", result)
        temp_paths.append(deterministic_out)
        if result["ok"] and _validate_candidate(deterministic_out, expected_pages):
            if output_path.exists():
                output_path.unlink()
            deterministic_out.replace(output_path)
            final_path = output_path

    if use_zopfli and qpdf:
        zopfli_out = _tmp_path("zopfli.pdf")
        env = os.environ.copy()
        env["QPDF_ZOPFLI"] = "enabled"
        result = _run_cmd(
            [
                "qpdf",
                "--object-streams=generate",
                "--compress-streams=y",
                "--recompress-flate",
                "--compression-level=9",
                "--deterministic-id",
                str(final_path),
                str(zopfli_out),
            ],
            timeout_seconds,
            env=env,
        )
        _record_step("qpdf_zopfli", result)
        temp_paths.append(zopfli_out)
        if result["ok"] and _validate_candidate(zopfli_out, expected_pages):
            zopfli_size = zopfli_out.stat().st_size
            zopfli_savings = max(size_bytes - zopfli_size, 0)
            if zopfli_savings >= min_savings_bytes and (
                zopfli_savings / size_bytes
            ) >= savings_threshold:
                if output_path.exists():
                    output_path.unlink()
                zopfli_out.replace(output_path)
                final_path = output_path

    for path in temp_paths:
        if path in {input_path, output_path}:
            continue
        try:
            if path.exists():
                path.unlink()
        except OSError:
            pass

    output_bytes = final_path.stat().st_size
    savings_bytes = max(size_bytes - output_bytes, 0)
    savings_percent = round((savings_bytes / size_bytes) * 100, 2) if size_bytes else 0.0

    result_payload = {
        "status": status,
        "method": method,
        "profile": profile,
        "original_bytes": size_bytes,
        "output_bytes": output_bytes,
        "savings_bytes": savings_bytes,
        "savings_percent": savings_percent,
        "steps": steps,
        "warnings": warnings,
        "image_metrics": image_metrics,
    }

    print(
        "Compression metrics:",
        {
            "profile": profile,
            "image_heavy": image_metrics.get("image_heavy"),
            "images": image_metrics.get("image_count"),
            "text_chars": image_metrics.get("text_chars"),
            "savings_pct": savings_percent,
        },
    )

    return final_path, result_payload


def _run_command(
    command: Sequence[str],
    timeout: int,
    error_prefix: str,
) -> subprocess.CompletedProcess[str]:
    """Run a shell command and raise a friendly error on timeout."""
    try:
        return subprocess.run(
            list(command),
            capture_output=True,
            text=True,
            check=False,
            timeout=timeout,
        )
    except subprocess.TimeoutExpired as error:
        raise RuntimeError(f"{error_prefix} timed out") from error


def _is_valid_pdf(path: Path) -> bool:
    """Basic PDF integrity check using qpdf when available plus pypdf parsing."""
    if not path.exists() or path.stat().st_size == 0:
        return False
    qpdf = shutil.which("qpdf")
    if qpdf:
        try:
            result = _run_command(
                [qpdf, "--check", str(path)],
                timeout=30,
                error_prefix="PDF validation",
            )
        except RuntimeError:
            return False
        if result.returncode != 0:
            return False
    try:
        reader = PdfReader(str(path))
        _ = len(reader.pages)
    except Exception:
        return False
    return True

def repair_pdf(input_path: Path, output_path: Path) -> Path:
    """Repair a PDF via staged strategies (qpdf -> mutool -> pypdf rewrite)."""
    reader: PdfReader | None = None
    try:
        reader = PdfReader(str(input_path))
        if reader.is_encrypted:
            raise ValueError("PDF is encrypted")
    except PdfReadError:
        # Corrupt PDFs should still get a chance through external repair tools.
        reader = None

    qpdf = shutil.which("qpdf")
    if qpdf:
        result = _run_command(
            [qpdf, "--linearize", str(input_path), str(output_path)],
            timeout=90,
            error_prefix="QPDF repair",
        )
        if result.returncode == 0 and _is_valid_pdf(output_path):
            return output_path

    mutool = shutil.which("mutool")
    if mutool:
        result = _run_command(
            [mutool, "clean", str(input_path), str(output_path)],
            timeout=90,
            error_prefix="MuPDF repair",
        )
        if result.returncode == 0 and _is_valid_pdf(output_path):
            return output_path

    if reader is None:
        raise RuntimeError("PDF is corrupted and could not be repaired")

    writer = PdfWriter()
    for page in reader.pages:
        writer.add_page(page)
    writer.add_metadata(reader.metadata or {})
    with output_path.open("wb") as handle:
        writer.write(handle)
    if not _is_valid_pdf(output_path):
        raise RuntimeError("Repair failed to produce a valid PDF")
    return output_path


def rotate_pdf(
    input_path: Path,
    output_path: Path,
    angle: int,
    pages: str | None,
    tolerant_ranges: bool = False,
) -> Path:
    """Rotate selected pages by the provided angle."""
    reader = _load_pdf(input_path)
    writer = PdfWriter()
    total_pages = len(reader.pages)
    target_pages = _resolve_page_selection(
        pages,
        total_pages,
        tolerant=tolerant_ranges,
    )
    for index, page in enumerate(reader.pages, start=1):
        if target_pages is None or index in target_pages:
            _rotate_page(page, angle)
        writer.add_page(page)
    _copy_metadata(writer, reader)
    with output_path.open("wb") as handle:
        writer.write(handle)
    return output_path


def watermark_pdf(
    input_path: Path,
    output_path: Path,
    text: str,
    pages: str | None,
    tolerant_ranges: bool = False,
) -> Path:
    """
    Apply a centered text watermark to selected pages of a PDF.
    
    Parameters:
        input_path (Path): Path to the source PDF.
        output_path (Path): Path where the watermarked PDF will be written.
        text (str): Watermark text to place centered on each target page.
        pages (str | None): Page selection expressed as a range string (e.g. "1-3,5"); pages are 1-based.
            If None, the watermark is applied to every page.
    
    Returns:
        Path: The same as `output_path` after the watermarked PDF has been written.
    """
    reader = _load_pdf(input_path)
    writer = PdfWriter()
    total_pages = len(reader.pages)
    target_pages = _resolve_page_selection(
        pages,
        total_pages,
        tolerant=tolerant_ranges,
    )
    for index, page in enumerate(reader.pages, start=1):
        if target_pages is None or index in target_pages:
            box = page.cropbox if hasattr(page, "cropbox") else page.mediabox
            width = float(box.upper_right[0] - box.lower_left[0])
            height = float(box.upper_right[1] - box.lower_left[1])

            def _draw(pdf: FPDF, width_mm: float, height_mm: float) -> None:
                """
                Render a large diagonal watermark from bottom-left to top-right.
                
                Calculates a font size from the smaller of the page width and height (clamped to 28-72), configures a unicode-capable font if required, sets a light gray color, and draws the text centered along a diagonal using rotation.
                
                Parameters:
                    pdf (FPDF): The FPDF instance representing the overlay page to draw on.
                    width_mm (float): Page width in millimeters.
                    height_mm (float): Page height in millimeters.
                """
                font_size = min(max(int(min(width_mm, height_mm) * 0.2), 28), 120)
                _set_overlay_font(pdf, text, font_size)
                pdf.set_text_color(165, 165, 165)
                angle = -math.degrees(math.atan2(height_mm, width_mm))
                center_x = width_mm / 2
                center_y = height_mm / 2
                with pdf.rotation(angle, x=center_x, y=center_y):
                    pdf.set_xy(0, center_y)
                    pdf.cell(width_mm, 10, text, align="C")

            overlay = _build_overlay_page(width, height, _draw)
            _merge_overlay_page(page, overlay, box)
        writer.add_page(page)
    _copy_metadata(writer, reader)
    with output_path.open("wb") as handle:
        writer.write(handle)
    return output_path


def page_numbers_pdf(
    input_path: Path,
    output_path: Path,
    start: int,
    pages: str | None,
    numbering_mode: str = "selectionIndex",
    tolerant_ranges: bool = False,
) -> Path:
    """
    Add sequential page numbers as footers to selected pages of a PDF.
    
    Parameters:
        input_path (Path): Path to the source PDF file.
        output_path (Path): Path where the resulting PDF will be written.
        start (int): Starting page number to apply to the first page (incremented per page).
        pages (str | None): Optional page selection string (e.g., "1-3,5") determining which pages receive numbers; if None, all pages are numbered.
    
    Returns:
        Path: The path to the written PDF file (same as output_path).
    """
    reader = _load_pdf(input_path)
    writer = PdfWriter()
    total_pages = len(reader.pages)
    target_pages = _resolve_page_selection(
        pages,
        total_pages,
        tolerant=tolerant_ranges,
    )
    numbering_mode_normalized = numbering_mode.strip() or "selectionIndex"
    if numbering_mode_normalized not in {"documentIndex", "selectionIndex"}:
        raise ValueError("Numbering mode must be documentIndex or selectionIndex")
    selection_counter = 0
    for index, page in enumerate(reader.pages, start=1):
        if target_pages is None or index in target_pages:
            box = page.cropbox if hasattr(page, "cropbox") else page.mediabox
            width = float(box.upper_right[0] - box.lower_left[0])
            height = float(box.upper_right[1] - box.lower_left[1])
            if numbering_mode_normalized == "selectionIndex":
                selection_counter += 1
                number = start + selection_counter - 1
            else:
                number = start + index - 1

            def _draw(pdf: FPDF, width_mm: float, height_mm: float) -> None:
                """
                Draws a centered numeric footer near the bottom edge of the overlay page.
                
                Positions and renders the page number (captured from the surrounding scope) as a footer using a font size chosen to fit the page: the size is proportional to the smaller page dimension and clamped to the range 8-16 points. The rendered text is centered with a 10 mm bottom margin and uses a muted gray color.
                
                Parameters:
                    pdf (FPDF): The FPDF instance used to draw on the overlay page.
                    width_mm (float): Page width in millimeters.
                    height_mm (float): Page height in millimeters.
                
                Returns:
                    None
                """
                font_size = min(max(int(min(width_mm, height_mm) * 0.04), 8), 16)
                _set_overlay_font(pdf, str(number), font_size)
                pdf.set_text_color(60, 60, 60)
                margin = 10
                pdf.set_xy(0, height_mm - margin)
                pdf.cell(width_mm, 6, str(number), align="C")

            overlay = _build_overlay_page(width, height, _draw)
            _merge_overlay_page(page, overlay, box)
        writer.add_page(page)
    _copy_metadata(writer, reader)
    with output_path.open("wb") as handle:
        writer.write(handle)
    return output_path


def crop_pdf(
    input_path: Path,
    output_path: Path,
    margins: str,
    pages: str | None,
    tolerant_ranges: bool = False,
) -> Path:
    """
    Crop selected pages of a PDF by the specified margins (measured in PDF points).
    
    Parameters:
        input_path (Path): Path to the source PDF.
        output_path (Path): Path where the cropped PDF will be written.
        margins (str): Margin specification in points; either a single numeric value applied to all sides or four comma-separated values in the order top,right,bottom,left.
        pages (str | None): Optional page selection string (e.g., "1-3,5") specifying which pages to crop; when None, all pages are processed.
    
    Returns:
        Path: The path to the written cropped PDF.
    
    Raises:
        ValueError: If the provided margins would remove an entire page or if margin parsing fails.
    """
    reader = _load_pdf(input_path)
    writer = PdfWriter()
    total_pages = len(reader.pages)
    top, right, bottom, left = _parse_margins(margins)
    if any(value < 0 for value in (top, right, bottom, left)):
        raise ValueError("Margins must be zero or positive")
    target_pages = _resolve_page_selection(
        pages,
        total_pages,
        tolerant=tolerant_ranges,
    )
    for index, page in enumerate(reader.pages, start=1):
        if target_pages is None or index in target_pages:
            lower_left_x = float(page.mediabox.lower_left[0]) + left
            lower_left_y = float(page.mediabox.lower_left[1]) + bottom
            upper_right_x = float(page.mediabox.upper_right[0]) - right
            upper_right_y = float(page.mediabox.upper_right[1]) - top
            if upper_right_x <= lower_left_x or upper_right_y <= lower_left_y:
                raise ValueError("Crop margins remove the entire page")
            page.mediabox.lower_left = (lower_left_x, lower_left_y)
            page.mediabox.upper_right = (upper_right_x, upper_right_y)
            page.cropbox.lower_left = (lower_left_x, lower_left_y)
            page.cropbox.upper_right = (upper_right_x, upper_right_y)
            page.trimbox.lower_left = (lower_left_x, lower_left_y)
            page.trimbox.upper_right = (upper_right_x, upper_right_y)
        writer.add_page(page)
    _copy_metadata(writer, reader)
    with output_path.open("wb") as handle:
        writer.write(handle)
    return output_path


def unlock_pdf(input_path: Path, output_path: Path, password: str = "") -> Path:
    """
    Remove password protection from the PDF at input_path and write the unlocked PDF to output_path.
    
    Preserves document pages and metadata. Raises ValueError if the PDF is encrypted and the provided password fails to decrypt it.
    
    Parameters:
        input_path (Path): Path to the source PDF (may be encrypted).
        output_path (Path): Path where the unlocked PDF will be written.
        password (str): Password to use for decryption.
    
    Returns:
        output_path (Path): The path to the written unlocked PDF.
    
    Raises:
        ValueError: If the PDF is encrypted and the provided password does not unlock it.
    """
    qpdf = shutil.which("qpdf")
    if qpdf:
        password_file: Path | None = None
        try:
            if not password:
                requires = subprocess.run(
                    [qpdf, "--requires-password", str(input_path)],
                    capture_output=True,
                    text=True,
                    check=False,
                    timeout=10,
                )
                if requires.returncode == 0:
                    raise ValueError("Password required to unlock this PDF")
                if requires.returncode == 2:
                    shutil.copyfile(input_path, output_path)
                    return output_path
            if password:
                with tempfile.NamedTemporaryFile(
                    mode="w",
                    encoding="utf-8",
                    delete=False,
                ) as handle:
                    handle.write(password)
                    password_file = Path(handle.name)
                password_file.chmod(0o600)
            cmd = [qpdf, "--decrypt", str(input_path), str(output_path)]
            if password_file is not None:
                cmd.insert(1, f"--password-file={password_file}")
            result = subprocess.run(
                cmd,
                capture_output=True,
                text=True,
                check=False,
                timeout=60,
            )
            if result.returncode == 0 and output_path.exists():
                return output_path
        except subprocess.TimeoutExpired as error:
            raise RuntimeError("Unlocking timed out") from error
        finally:
            if password_file is not None:
                password_file.unlink(missing_ok=True)

    reader = PdfReader(str(input_path))
    if reader.is_encrypted:
        result = reader.decrypt(password or "")
        if result == 0:
            raise ValueError("Password required to unlock this PDF")
    writer = PdfWriter()
    for page in reader.pages:
        writer.add_page(page)
    writer.add_metadata(reader.metadata or {})
    with output_path.open("wb") as handle:
        writer.write(handle)
    return output_path


def protect_pdf(input_path: Path, output_path: Path, password: str) -> Path:
    """
    Encrypts an unencrypted PDF file with the specified password.
    
    Sets the same value as both the user and owner password and enables 128-bit encryption; preserves the input PDF's metadata.
    
    Parameters:
        input_path (Path): Path to the source PDF file (must be unencrypted).
        output_path (Path): Path where the encrypted PDF will be written.
        password (str): Password to apply as both the user and owner password.
    
    Returns:
        Path: The path to the written encrypted PDF (the provided output_path).
    
    Raises:
        ValueError: If the input PDF is already encrypted.
    """
    reader = PdfReader(str(input_path))
    if reader.is_encrypted:
        raise ValueError("PDF is already encrypted")
    writer = PdfWriter()
    for page in reader.pages:
        writer.add_page(page)
    writer.add_metadata(reader.metadata or {})
    writer.encrypt(user_password=password, owner_password=password, use_128bit=True)
    with output_path.open("wb") as handle:
        writer.write(handle)
    return output_path


def redact_pdf(
    input_path: Path,
    output_path: Path,
    text: str,
    pages: str | None,
    case_sensitive: bool = False,
    whole_word: bool = False,
    use_regex: bool = False,
    ocr_assist: bool = False,
    tolerant_ranges: bool = False,
) -> Path:
    """
    Redact all occurrences of a given text in a PDF, optionally restricted to specific pages.
    
    Searches each targeted page for exact occurrences of `text`, adds black redact annotations over matches, applies the redactions, and writes the modified PDF to `output_path`.
    
    Parameters:
        input_path (Path): Path to the source PDF.
        output_path (Path): Path where the redacted PDF will be written.
        text (str): Text to search for and redact; matches are searched as exact occurrences.
        pages (str | None): Optional page selection string (e.g., "1-3,5") restricting which pages to process; if `None`, all pages are searched.
    
    When `ocr_assist` is enabled and searchable text is not present, OCR bounding boxes
    are used for redaction candidates. If OCR locations cannot be derived reliably but a
    text match is detected, the function falls back to full-page redaction and logs a warning.

    Returns:
        Path: The `output_path` of the saved redacted PDF.
    """
    query = text.strip()
    if not query:
        raise ValueError("Text to redact is required")

    regex: re.Pattern[str] | None = None
    if use_regex:
        try:
            regex = re.compile(query, 0 if case_sensitive else re.IGNORECASE)
        except re.error as error:
            raise ValueError(f"Invalid regex pattern: {query!r}") from error

    with fitz.open(str(input_path)) as document:
        _assert_fitz_unencrypted(document)
        total_pages = document.page_count
        target_pages = _resolve_page_selection(
            pages,
            total_pages,
            tolerant=tolerant_ranges,
        )
        for index in range(total_pages):
            page_number = index + 1
            if target_pages is not None and page_number not in target_pages:
                continue
            page = document.load_page(index)
            rectangles: list[fitz.Rect] = []

            if not use_regex and not whole_word and not case_sensitive:
                rectangles = page.search_for(query)
            else:
                words = page.get_text("words") or []
                for item in words:
                    if len(item) < 5:
                        continue
                    rect = fitz.Rect(item[0], item[1], item[2], item[3])
                    token = str(item[4])
                    candidate = token if case_sensitive else token.lower()
                    target = query if case_sensitive else query.lower()
                    if use_regex:
                        if regex is not None and regex.search(token):
                            rectangles.append(rect)
                    elif whole_word:
                        if candidate == target:
                            rectangles.append(rect)
                    else:
                        if target in candidate:
                            rectangles.append(rect)

            if not rectangles and ocr_assist and pytesseract is not None:
                image = _render_page_image(page, OCR_DPI)
                matched_text = False
                try:
                    ocr_data = pytesseract.image_to_data(
                        image,
                        lang=DEFAULT_OCR_LANG,
                        output_type=getattr(pytesseract, "Output", object()).DICT,
                    )
                    texts = ocr_data.get("text", [])
                    lefts = ocr_data.get("left", [])
                    tops = ocr_data.get("top", [])
                    widths = ocr_data.get("width", [])
                    heights = ocr_data.get("height", [])
                    if (
                        texts
                        and len(texts) == len(lefts) == len(tops) == len(widths) == len(heights)
                    ):
                        img_w, img_h = image.size
                        page_rect = page.rect
                        scale_x = page_rect.width / max(1, img_w)
                        scale_y = page_rect.height / max(1, img_h)
                        for token, left, top, width, height in zip(
                            texts, lefts, tops, widths, heights
                        ):
                            token_text = str(token).strip()
                            if not token_text:
                                continue
                            if use_regex and regex is not None:
                                token_match = regex.search(token_text) is not None
                            elif whole_word:
                                if case_sensitive:
                                    token_match = token_text == query
                                else:
                                    token_match = token_text.lower() == query.lower()
                            else:
                                if case_sensitive:
                                    token_match = query in token_text
                                else:
                                    token_match = query.lower() in token_text.lower()
                            if not token_match:
                                continue
                            matched_text = True
                            x0 = page_rect.x0 + float(left) * scale_x
                            y0 = page_rect.y0 + float(top) * scale_y
                            x1 = x0 + float(width) * scale_x
                            y1 = y0 + float(height) * scale_y
                            if x1 > x0 and y1 > y0:
                                rectangles.append(fitz.Rect(x0, y0, x1, y1))
                except Exception:
                    # Continue with text-level fallback below.
                    pass

                if not rectangles:
                    page_text = _ocr_image(image, DEFAULT_OCR_LANG)
                    if use_regex and regex is not None:
                        matched_text = regex.search(page_text) is not None
                    elif whole_word:
                        words = re.findall(r"\b\w+\b", page_text)
                        if case_sensitive:
                            matched_text = query in words
                        else:
                            matched_text = query.lower() in [word.lower() for word in words]
                    else:
                        if case_sensitive:
                            matched_text = query in page_text
                        else:
                            matched_text = query.lower() in page_text.lower()
                    if matched_text:
                        LOGGER.warning(
                            "OCR assist fell back to full-page redaction on page %s for query %r",
                            page_number,
                            query,
                        )
                        rectangles = [page.rect]

            if not rectangles:
                continue
            for rect in rectangles:
                page.add_redact_annot(rect, fill=(0, 0, 0))
            page.apply_redactions()
        document.save(
            str(output_path),
            deflate=True,
            garbage=4,
            clean=True,
            incremental=False,
        )
    return output_path


def compare_pdfs(
    first_path: Path,
    second_path: Path,
    output_path: Path,
    include_visual_diff: bool = True,
) -> Path:
    """
    Produce a plain-text comparison report summarizing page counts and per-page text differences between two PDFs.
    
    The report includes the input filenames, their page counts, and entries for pages that are missing from one file or whose extracted text differs; if no differences are found the report records that fact. The report is written to output_path using UTF-8.
    
    Parameters:
        first_path (Path): Path to the first PDF file (File A).
        second_path (Path): Path to the second PDF file (File B).
        output_path (Path): Destination path for the generated text report.
    
    Returns:
        Path: The path to the written report (output_path).
    """
    def _regions_from_diff(diff_image: Image.Image, tile_size: int = 48) -> list[tuple[int, int, int, int]]:
        """Return merged changed regions from a grayscale diff image."""
        binary = diff_image.point(lambda px: 255 if px > 12 else 0, mode="L")
        width, height = binary.size
        if width <= 0 or height <= 0:
            return []
        tiles: list[tuple[int, int, int, int]] = []
        for y in range(0, height, tile_size):
            for x in range(0, width, tile_size):
                x2 = min(width, x + tile_size)
                y2 = min(height, y + tile_size)
                tile = binary.crop((x, y, x2, y2))
                histogram = tile.histogram()
                changed = histogram[255] if len(histogram) > 255 else 0
                tile_pixels = max(1, tile.width * tile.height)
                if changed / tile_pixels >= 0.03:
                    tiles.append((x, y, x2, y2))
        if not tiles:
            return []

        merged: list[tuple[int, int, int, int]] = []
        for tile in tiles:
            tx0, ty0, tx1, ty1 = tile
            merged_any = False
            for index, region in enumerate(merged):
                rx0, ry0, rx1, ry1 = region
                if tx0 <= rx1 + tile_size and tx1 >= rx0 - tile_size and ty0 <= ry1 + tile_size and ty1 >= ry0 - tile_size:
                    merged[index] = (
                        min(rx0, tx0),
                        min(ry0, ty0),
                        max(rx1, tx1),
                        max(ry1, ty1),
                    )
                    merged_any = True
                    break
            if not merged_any:
                merged.append(tile)

        merged.sort(key=lambda item: (item[2] - item[0]) * (item[3] - item[1]), reverse=True)
        return merged[:8]

    reader_a = _load_pdf(first_path)
    reader_b = _load_pdf(second_path)
    pages_a = len(reader_a.pages)
    pages_b = len(reader_b.pages)
    lines = [
        "ZenPDF comparison report",
        f"File A: {first_path.name}",
        f"File B: {second_path.name}",
        f"Pages: {pages_a} vs {pages_b}",
        "",
    ]
    differences: List[str] = []
    visual_notes: List[str] = []
    if pages_a != pages_b:
        differences.append("Page counts differ.")
    for index in range(max(pages_a, pages_b)):
        if index >= pages_a:
            differences.append(f"Page {index + 1}: missing from file A")
            continue
        if index >= pages_b:
            differences.append(f"Page {index + 1}: missing from file B")
            continue
        text_a = (reader_a.pages[index].extract_text() or "").strip()
        text_b = (reader_b.pages[index].extract_text() or "").strip()
        if text_a != text_b:
            differences.append(f"Page {index + 1}: text differs")

    if include_visual_diff:
        try:
            with fitz.open(str(first_path)) as doc_a, fitz.open(str(second_path)) as doc_b:
                max_pages = max(doc_a.page_count, doc_b.page_count)
                matrix = fitz.Matrix(1, 1)
                for page_index in range(max_pages):
                    if page_index >= doc_a.page_count or page_index >= doc_b.page_count:
                        continue
                    pix_a = doc_a.load_page(page_index).get_pixmap(
                        matrix=matrix, colorspace=fitz.csGRAY, alpha=False
                    )
                    pix_b = doc_b.load_page(page_index).get_pixmap(
                        matrix=matrix, colorspace=fitz.csGRAY, alpha=False
                    )
                    image_a = Image.open(BytesIO(pix_a.tobytes("png"))).convert("L")
                    image_b = Image.open(BytesIO(pix_b.tobytes("png"))).convert("L")
                    if image_a.size != image_b.size:
                        visual_notes.append(
                            f"Page {page_index + 1}: visual size differs ({image_a.size} vs {image_b.size})"
                        )
                        continue
                    diff = ImageChops.difference(image_a, image_b)
                    histogram = diff.histogram()
                    changed = sum(histogram[1:])
                    total_pixels = image_a.size[0] * image_a.size[1]
                    if total_pixels == 0:
                        continue
                    percent = (changed / total_pixels) * 100
                    if percent > 0.1:
                        regions = _regions_from_diff(diff)
                        visual_notes.append(
                            f"Page {page_index + 1}: visual delta {percent:.2f}% ({len(regions)} region(s))"
                        )
                        for region_index, (x0, y0, x1, y1) in enumerate(regions, start=1):
                            visual_notes.append(
                                f"Page {page_index + 1}: region {region_index} bbox=({x0},{y0})-({x1},{y1})"
                            )
        except Exception as error:
            visual_notes.append(f"Visual diff unavailable: {error}")
    if not differences:
        lines.append("No text differences detected.")
    else:
        lines.extend(differences)
    if include_visual_diff:
        lines.append("")
        lines.append("Visual comparison")
        if visual_notes:
            lines.extend(visual_notes)
        else:
            lines.append("No visual differences detected.")
    output_path.write_text("\n".join(lines), encoding="utf-8")
    return output_path


def image_to_pdf(inputs: Sequence[Path], output_path: Path) -> Path:
    """
    Combine one or more image files into a single PDF file.
    
    Parameters:
        inputs (Sequence[Path]): Paths to image files to include, in order.
        output_path (Path): Destination path for the generated PDF.
    
    Returns:
        Path: The path to the written PDF (same object as `output_path`).
    
    Raises:
        ValueError: If image rendering to PDF fails.
    """
    pdf_bytes = img2pdf.convert([str(path) for path in inputs])
    if pdf_bytes is None:
        raise ValueError("Failed to render images to PDF")
    output_path.write_bytes(pdf_bytes)
    return output_path


def pdf_to_jpg(input_path: Path, output_dir: Path, dpi: int = 150) -> List[Path]:
    """Render each PDF page to a JPG image."""
    with fitz.open(str(input_path)) as document:
        _assert_fitz_unencrypted(document)
        scale = dpi / 72
        matrix = fitz.Matrix(scale, scale)
        outputs: List[Path] = []
        stem = input_path.stem
        if "_" in stem:
            prefix, remainder = stem.split("_", 1)
            if prefix.isdigit() and len(prefix) == 2:
                stem = remainder
        for index in range(document.page_count):
            page = document.load_page(index)
            render = getattr(page, "get_pixmap", None)
            if not callable(render):
                raise ValueError("PDF renderer unavailable")
            pix = render(matrix=matrix)
            output_path = output_dir / f"{stem}_{index + 1}.jpg"
            saver = getattr(pix, "save", None)
            if not callable(saver):
                raise ValueError("Rendered page cannot be saved")
            saver(str(output_path))
            outputs.append(output_path)
        return outputs


def zip_outputs(outputs: Iterable[Path], zip_path: Path) -> Path:
    """Zip multiple output files into a single archive."""
    with zipfile.ZipFile(zip_path, "w", compression=zipfile.ZIP_DEFLATED) as archive:
        for item in outputs:
            archive.write(item, arcname=item.name)
    return zip_path


MAX_WEB_BYTES = 2 * 1024 * 1024
UNICODE_FONT_PATHS = (
    Path(__file__).resolve().parent / "assets" / "DejaVuSans.ttf",
    Path(__file__).resolve().parent / "assets" / "NotoSans-Regular.ttf",
    Path("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"),
    Path("/usr/share/fonts/DejaVuSans.ttf"),
)


def _resolve_unicode_font_path() -> Path | None:
    """
    Locate a Unicode-compatible TrueType font file if one is available.
    
    Checks the ZENPDF_TTF_PATH environment variable first, then falls back to known candidate paths.
    
    Returns:
        Path | None: Path to the font file if found, `None` otherwise.
    """
    env_path = os.getenv("ZENPDF_TTF_PATH")
    if env_path:
        candidate = Path(env_path)
        if candidate.is_file():
            return candidate
    for candidate in UNICODE_FONT_PATHS:
        if candidate.is_file():
            return candidate
    return None


def _set_overlay_font(pdf: FPDF, text: str, size: int) -> None:
    """
    Selects and configures an appropriate font on the given FPDF instance for rendering overlay text.
    
    Attempts to load a Unicode-capable DejaVu Sans from the environment or known paths; if unavailable, verifies whether the provided text can be encoded in Latin-1 and falls back to Helvetica. If the text requires Unicode and no Unicode font is available, raises RuntimeError.
    
    Parameters:
        pdf (FPDF): The FPDF instance to configure.
        text (str): Sample text to test whether a Unicode font is required.
        size (int): Font size to set on the PDF.
    
    Raises:
        RuntimeError: If the text contains characters that require a Unicode font but no Unicode font path is available.
    """
    font_path = _resolve_unicode_font_path()
    if font_path:
        pdf.add_font("DejaVuSans", fname=str(font_path))
        pdf.set_font("DejaVuSans", size=size)
        return
    try:
        text.encode("latin-1")
    except UnicodeEncodeError as error:
        raise RuntimeError(
            "Unicode font unavailable. Set ZENPDF_TTF_PATH to a Unicode TTF (e.g. DejaVuSans.ttf or NotoSans-Regular.ttf)."
        ) from error
    pdf.set_font("Helvetica", size=size)


class _HTMLTextExtractor(HTMLParser):
    """Minimal HTML to text extractor."""

    def __init__(self) -> None:
        """Initialize the parser."""
        super().__init__()
        self._parts: List[str] = []

    def handle_data(self, data: str) -> None:
        """Capture text content from HTML."""
        text = data.strip()
        if text:
            self._parts.append(text)

    def text(self) -> str:
        """Return the collected text content."""
        return "\n".join(self._parts)


def html_to_pdf(html: str, output_path: Path) -> Path:
    """Render basic HTML text into a PDF."""
    parser = _HTMLTextExtractor()
    parser.feed(html)
    text = parser.text() or "(no content)"

    pdf = FPDF(orientation="P", unit="mm", format="A4")
    pdf.set_margins(15, 15, 15)
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()
    font_path = _resolve_unicode_font_path()
    if font_path:
        pdf.add_font("DejaVuSans", fname=str(font_path))
        pdf.set_font("DejaVuSans", size=12)
    else:
        try:
            text.encode("latin-1")
        except UnicodeEncodeError as error:
            raise RuntimeError(
                "Unicode font unavailable. Set ZENPDF_TTF_PATH to a Unicode TTF (e.g. DejaVuSans.ttf or NotoSans-Regular.ttf)."
            ) from error
        pdf.set_font("Helvetica", size=12)
    max_width = pdf.w - pdf.l_margin - pdf.r_margin
    for line in text.splitlines():
        if line.strip():
            pdf.multi_cell(max_width, 6, line)
        else:
            pdf.ln(5)
    pdf.output(str(output_path))
    return output_path


def _resolve_browser_binary() -> str | None:
    """Return a browser executable that supports --print-to-pdf."""
    env_value = os.getenv("ZENPDF_BROWSER_PATH")
    if env_value:
        candidate = shutil.which(env_value) if not Path(env_value).is_absolute() else env_value
        if candidate:
            return str(candidate)
    for candidate in (
        "chromium",
        "chromium-browser",
        "google-chrome",
        "google-chrome-stable",
        "chrome",
    ):
        resolved = shutil.which(candidate)
        if resolved:
            return resolved
    return None


def _is_dev_mode() -> bool:
    """Return True for local/dev runtime toggles."""
    return os.getenv("ZENPDF_DEV_MODE") == "1" or os.getenv("NODE_ENV") == "development"


def _validate_browser_request_url(request_url: str) -> None:
    """Validate browser network targets against SSRF policy."""
    parsed = urlparse(request_url)
    if parsed.scheme in {"data", "blob", "about"}:
        return
    if parsed.scheme not in {"http", "https"} or not parsed.hostname:
        raise ValueError(f"Unsupported browser request URL: {request_url!r}")
    _resolve_public_ip(parsed.hostname)


def _browser_url_to_pdf(url: str, output_path: Path) -> Path:
    """Render a URL into PDF using Playwright with strict network policy."""
    if sync_playwright is None:
        raise RuntimeError("Browser mode requires Playwright in worker runtime")

    browser = _resolve_browser_binary()
    blocked_requests: list[str] = []
    output_path.parent.mkdir(parents=True, exist_ok=True)
    allow_insecure = (
        os.getenv("ZENPDF_WEB_ALLOW_INSECURE_SSL") == "1"
        and _is_dev_mode()
    )

    launch_args = [
        "--disable-gpu",
        "--disable-dev-shm-usage",
        "--no-first-run",
        "--no-default-browser-check",
        "--disable-features=TranslateUI",
    ]

    try:
        with sync_playwright() as playwright:
            launch_kwargs: dict[str, Any] = {
                "headless": True,
                "args": launch_args,
            }
            if browser:
                launch_kwargs["executable_path"] = browser
            browser_instance = playwright.chromium.launch(**launch_kwargs)
            context = browser_instance.new_context(ignore_https_errors=allow_insecure)
            page = context.new_page()

            def _route_handler(route, request) -> None:  # type: ignore[no-untyped-def]
                request_url = request.url
                try:
                    _validate_browser_request_url(request_url)
                except Exception:
                    blocked_requests.append(request_url)
                    route.abort()
                    return
                route.continue_()

            page.route("**/*", _route_handler)
            response = page.goto(url, wait_until="domcontentloaded", timeout=60_000)
            if response and response.status >= 400:
                raise RuntimeError(f"Browser render failed with status {response.status}")
            page.wait_for_load_state("networkidle", timeout=30_000)
            if blocked_requests:
                blocked_sample = ", ".join(blocked_requests[:3])
                raise ValueError(
                    f"Blocked private or unsupported browser request(s): {blocked_sample}"
                )
            page.emulate_media(media="screen")
            page.pdf(
                path=str(output_path),
                print_background=True,
                prefer_css_page_size=True,
            )
            context.close()
            browser_instance.close()
    except PlaywrightError as error:
        raise RuntimeError(f"Browser PDF render failed: {error}") from error

    if not output_path.exists() or output_path.stat().st_size == 0:
        raise RuntimeError("Browser PDF render produced no output")
    return output_path


def web_to_pdf(url: str, output_path: Path, render_mode: str = "browser") -> Path:
    """Fetch a URL and convert it to PDF via browser render or text extraction."""
    parsed = urlparse(url)
    if parsed.scheme not in {"http", "https"} or not parsed.hostname:
        raise ValueError("Only http/https URLs are supported")
    render_mode_normalized = render_mode.strip().lower() or "browser"
    if render_mode_normalized not in {"browser", "text"}:
        raise ValueError("Render mode must be browser or text")

    _ = _resolve_public_ip(parsed.hostname)

    # Preflight to prevent redirect-based bypass and unsupported responses.
    allow_hostname_fallback = parsed.scheme == "https" and (
        os.getenv("ZENPDF_WEB_ALLOW_HOSTNAME_FALLBACK") == "1"
        and _is_dev_mode()
    )

    try:
        with requests.Session() as preflight:
            response = preflight.get(url, allow_redirects=False, timeout=20, stream=True)
            with response:
                if 300 <= response.status_code < 400:
                    raise ValueError("Redirects are not allowed")
                response.raise_for_status()
    except requests.exceptions.SSLError:
        if not allow_hostname_fallback:
            raise

    if render_mode_normalized == "browser":
        return _browser_url_to_pdf(parsed.geturl(), output_path)

    target_ip = _resolve_public_ip(parsed.hostname)
    is_ipv6 = isinstance(ipaddress.ip_address(target_ip), ipaddress.IPv6Address)
    host = f"[{target_ip}]" if is_ipv6 else target_ip
    netloc = f"{host}:{parsed.port}" if parsed.port is not None else host
    try:
        hostname_ip = ipaddress.ip_address(parsed.hostname)
        hostname_is_ipv6 = isinstance(hostname_ip, ipaddress.IPv6Address)
    except ValueError:
        hostname_is_ipv6 = False

    host_header = (
        f"[{parsed.hostname}]" if hostname_is_ipv6 else parsed.hostname
    )
    if parsed.port is not None:
        host_header = f"{host_header}:{parsed.port}"
    target_url = parsed._replace(netloc=netloc).geturl()

    def _fetch_html(
        session: requests.Session,
        fetch_url: str,
        header: str | None,
        verify_ssl: bool,
    ) -> tuple[bytearray, str]:
        body = bytearray()
        headers = {"Host": header} if header else None
        with session.get(
            fetch_url,
            timeout=20,
            allow_redirects=False,
            stream=True,
            headers=headers,
            verify=verify_ssl,
        ) as response:
            if 300 <= response.status_code < 400:
                raise ValueError("Redirects are not allowed")
            response.raise_for_status()
            for chunk in response.iter_content(chunk_size=64 * 1024):
                if not chunk:
                    continue
                body.extend(chunk)
                if len(body) > MAX_WEB_BYTES:
                    raise ValueError("Web response too large")
            encoding = response.encoding or "utf-8"
        return body, encoding

    body: bytearray
    encoding: str
    allow_insecure = (
        os.getenv("ZENPDF_WEB_ALLOW_INSECURE_SSL") == "1"
        and _is_dev_mode()
    )

    with requests.Session() as session:
        if parsed.scheme == "https":
            session.mount("https://", HostHeaderSSLAdapter())

        try:
            body, encoding = _fetch_html(
                session,
                target_url,
                host_header,
                not allow_insecure,
            )
        except requests.exceptions.SSLError:
            allow_fallback = parsed.scheme == "https" and (
                os.getenv("ZENPDF_WEB_ALLOW_HOSTNAME_FALLBACK") == "1"
                and _is_dev_mode()
            )
            if not allow_fallback:
                raise
            # Re-validate hostname before falling back to hostname-based HTTPS.
            _resolve_public_ip(parsed.hostname)
            with requests.Session() as fallback_session:
                body, encoding = _fetch_html(
                    fallback_session,
                    parsed.geturl(),
                    None,
                    not allow_insecure,
                )

    html = body.decode(encoding, errors="replace")
    return html_to_pdf(html, output_path)


def office_to_pdf(input_path: Path, output_dir: Path) -> Path:
    """Convert an Office document to PDF using LibreOffice."""
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        raise RuntimeError("LibreOffice is required for Office to PDF conversion")

    output_dir.mkdir(parents=True, exist_ok=True)

    size_mb = max(1, math.ceil(input_path.stat().st_size / (1024 * 1024)))
    timeout_base = int(os.getenv("ZENPDF_OFFICE_TIMEOUT_BASE_SECONDS", "120"))
    timeout_per_mb = int(os.getenv("ZENPDF_OFFICE_TIMEOUT_PER_MB_SECONDS", "4"))
    timeout_max = int(os.getenv("ZENPDF_OFFICE_TIMEOUT_MAX_SECONDS", "480"))
    timeout_seconds = min(timeout_max, timeout_base + size_mb * timeout_per_mb)

    output_path = output_dir / f"{input_path.stem}.pdf"
    command = [
        soffice,
        "--headless",
        "--nologo",
        "--norestore",
        "--nolockcheck",
        "--nodefault",
        "--convert-to",
        "pdf",
        "--outdir",
        str(output_dir),
        str(input_path),
    ]

    last_error_detail = ""
    attempts = [
        ("primary", timeout_seconds, {}),
        (
            "retry-safe",
            min(timeout_max * 2, max(timeout_seconds + 60, int(timeout_seconds * 1.8))),
            {"SAL_USE_VCLPLUGIN": "svp", "SAL_DISABLE_OPENCL": "1"},
        ),
    ]
    for index, (attempt_name, attempt_timeout, env_patch) in enumerate(attempts, start=1):
        if output_path.exists():
            output_path.unlink(missing_ok=True)
        env = os.environ.copy()
        env.update(env_patch)
        try:
            result = subprocess.run(
                command,
                capture_output=True,
                text=True,
                check=False,
                timeout=attempt_timeout,
                env=env,
            )
        except subprocess.TimeoutExpired:
            last_error_detail = f"{attempt_name} timed out after {attempt_timeout}s"
            if index < len(attempts):
                LOGGER.warning("Office conversion attempt %s timed out; retrying", attempt_name)
                continue
            raise RuntimeError("Office conversion timed out. Try a smaller file.")

        detail = (result.stderr or result.stdout or "").strip()
        if result.returncode == 0 and output_path.exists():
            return output_path
        last_error_detail = detail or f"{attempt_name} failed with exit code {result.returncode}"
        detail_lower = last_error_detail.lower()
        if "password" in detail_lower or "encrypted" in detail_lower:
            raise ValueError("Office file appears password-protected or encrypted")
        if "font" in detail_lower:
            raise RuntimeError(
                "Office conversion failed due to missing fonts. Install font packs (e.g. DejaVu/Noto/Liberation) and retry."
            )
        if index < len(attempts):
            LOGGER.warning(
                "Office conversion attempt %s failed (%s); retrying",
                attempt_name,
                last_error_detail,
            )
            continue

    raise RuntimeError(
        "Office conversion failed. Please verify the source file format."
        + (f" Details: {last_error_detail}" if last_error_detail else "")
    )


def _ensure_extension(input_path: Path, allowed_extensions: set[str], label: str) -> None:
    """Validate input extension for Office conversion routes."""
    extension = input_path.suffix.lower()
    if extension not in allowed_extensions:
        supported = ", ".join(sorted(allowed_extensions))
        raise ValueError(f"{label} expects one of: {supported}")


def word_to_pdf(input_path: Path, output_dir: Path) -> Path:
    """Convert a Word file to PDF using LibreOffice."""
    _ensure_extension(input_path, {".doc", ".docx"}, "Word to PDF")
    return office_to_pdf(input_path, output_dir)


def powerpoint_to_pdf(input_path: Path, output_dir: Path) -> Path:
    """Convert a PowerPoint file to PDF using LibreOffice."""
    _ensure_extension(input_path, {".ppt", ".pptx"}, "PowerPoint to PDF")
    return office_to_pdf(input_path, output_dir)


def excel_to_pdf(input_path: Path, output_dir: Path) -> Path:
    """Convert an Excel file to PDF using LibreOffice."""
    _ensure_extension(input_path, {".xls", ".xlsx"}, "Excel to PDF")
    return office_to_pdf(input_path, output_dir)


def pdf_to_powerpoint(
    input_path: Path,
    output_path: Path,
    dpi: int = 150,
    mode: str = "visual",
) -> Path:
    """Convert PDF pages to PPTX in visual-fidelity or editable best-effort mode."""
    if Presentation is None:
        raise RuntimeError("python-pptx is required for PDF to PowerPoint")
    normalized_mode = (mode or "visual").strip().lower()
    if normalized_mode not in {"visual", "editable"}:
        raise ValueError("Mode must be visual or editable")
    if normalized_mode == "editable" and (Emu is None or Pt is None):
        raise RuntimeError("Editable PPT mode requires python-pptx utilities")

    def _fit_image_to_slide(
        image_width: int,
        image_height: int,
        target_width: int,
        target_height: int,
    ) -> tuple[int, int, int, int]:
        """Return centered (left, top, width, height) preserving image aspect ratio."""
        image_ratio = image_width / image_height
        target_ratio = target_width / target_height
        if image_ratio > target_ratio:
            width = target_width
            height = max(1, int(round(width / image_ratio)))
            left = 0
            top = max(0, (target_height - height) // 2)
        else:
            height = target_height
            width = max(1, int(round(height * image_ratio)))
            top = 0
            left = max(0, (target_width - width) // 2)
        return left, top, width, height

    with fitz.open(str(input_path)) as document:
        _assert_fitz_unencrypted(document)
        if document.page_count == 0:
            raise ValueError("PDF has no pages")
        presentation = Presentation()
        layouts = presentation.slide_layouts
        if len(layouts) == 0:
            raise RuntimeError("No slide layouts are available for PowerPoint export")
        blank_layout = layouts[6] if len(layouts) > 6 else layouts[len(layouts) - 1]
        scale = dpi / 72
        matrix = fitz.Matrix(scale, scale)
        first_page = document.load_page(0)
        # PDF points (1/72 inch) -> EMU (914400/72 = 12700)
        presentation.slide_width = max(1, int(round(first_page.rect.width * 12700)))
        presentation.slide_height = max(1, int(round(first_page.rect.height * 12700)))
        slide_width = int(presentation.slide_width)
        slide_height = int(presentation.slide_height)
        for index in range(document.page_count):
            page = document.load_page(index)
            slide = presentation.slides.add_slide(blank_layout)
            if normalized_mode == "visual":
                pix = page.get_pixmap(matrix=matrix, alpha=False)
                image_bytes = pix.tobytes("png")
                left, top, width, height = _fit_image_to_slide(
                    pix.width,
                    pix.height,
                    slide_width,
                    slide_height,
                )
                slide.shapes.add_picture(
                    BytesIO(image_bytes),
                    left,
                    top,
                    width=width,
                    height=height,
                )
                continue

            # Editable mode: place text boxes and extracted images as a best-effort mapping.
            page_dict = page.get_text("dict")
            blocks = page_dict.get("blocks", [])
            for block in blocks:
                block_type = int(block.get("type", -1))
                x0 = max(0.0, float(block.get("bbox", [0, 0, 0, 0])[0]))
                y0 = max(0.0, float(block.get("bbox", [0, 0, 0, 0])[1]))
                x1 = max(x0 + 1.0, float(block.get("bbox", [0, 0, 1, 1])[2]))
                y1 = max(y0 + 1.0, float(block.get("bbox", [0, 0, 1, 1])[3]))
                left = Emu(int(round(x0 * 12700)))
                top = Emu(int(round(y0 * 12700)))
                width = Emu(int(round((x1 - x0) * 12700)))
                height = Emu(int(round((y1 - y0) * 12700)))

                if block_type == 0:
                    lines: list[str] = []
                    for line in block.get("lines", []):
                        spans = line.get("spans", [])
                        line_text = "".join(str(span.get("text", "")) for span in spans).strip()
                        if line_text:
                            lines.append(line_text)
                    text = "\n".join(lines).strip()
                    if not text:
                        continue
                    textbox = slide.shapes.add_textbox(left, top, width, height)
                    frame = textbox.text_frame
                    frame.clear()
                    frame.word_wrap = True
                    frame.text = text
                    for paragraph in frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(11)
                    continue

                if block_type == 1:
                    image_bytes = block.get("image")
                    if isinstance(image_bytes, (bytes, bytearray)):
                        slide.shapes.add_picture(BytesIO(bytes(image_bytes)), left, top, width, height)
    presentation.save(str(output_path))
    return output_path


def scan_to_pdf(inputs: Sequence[Path], output_path: Path) -> Path:
    """Create a scan PDF from one or more images."""
    return image_to_pdf(inputs, output_path)


def sign_pdf(
    input_path: Path,
    output_path: Path,
    text: str,
    pages: str | None = None,
    x: float = 36.0,
    y: float = 36.0,
    font_size: float = 18.0,
    anchor: str = "custom",
    mode: str = "visual",
    signature_image_path: Path | None = None,
    pkcs12_path: Path | None = None,
    pkcs12_password: str | None = None,
    tolerant_ranges: bool = False,
) -> Path:
    """Sign a PDF using either visual stamp mode or cryptographic signature mode."""

    def _resolve_signature_rect(
        page_rect: fitz.Rect,
        anchor_mode: str,
        box_width: float,
        box_height: float,
        x_pos: float,
        y_pos: float,
    ) -> fitz.Rect:
        margin = 24.0
        target_x = x_pos
        target_y = y_pos
        if anchor_mode == "bottom-right":
            target_x = page_rect.x1 - box_width - margin
            target_y = page_rect.y1 - box_height - margin
        elif anchor_mode == "bottom-left":
            target_x = page_rect.x0 + margin
            target_y = page_rect.y1 - box_height - margin
        elif anchor_mode == "top-right":
            target_x = page_rect.x1 - box_width - margin
            target_y = page_rect.y0 + margin
        elif anchor_mode == "top-left":
            target_x = page_rect.x0 + margin
            target_y = page_rect.y0 + margin

        target_x = min(max(target_x, page_rect.x0 + 1), max(page_rect.x0 + 1, page_rect.x1 - 41))
        target_y = min(max(target_y, page_rect.y0 + 1), max(page_rect.y0 + 1, page_rect.y1 - 25))
        if target_x + box_width > page_rect.x1:
            box_width = max(40.0, page_rect.x1 - target_x - 1)
        if target_y + box_height > page_rect.y1:
            box_height = max(24.0, page_rect.y1 - target_y - 1)
        return fitz.Rect(target_x, target_y, target_x + box_width, target_y + box_height)

    normalized_mode = (mode or "visual").strip().lower()
    if normalized_mode not in {"visual", "cryptographic"}:
        raise ValueError("Sign mode must be visual or cryptographic")

    if normalized_mode == "cryptographic":
        if signers is None or IncrementalPdfFileWriter is None:
            raise RuntimeError("Cryptographic signing requires pyHanko")
        if pkcs12_path is None:
            raise ValueError("PKCS#12 certificate (.p12/.pfx) is required for cryptographic signing")
        if not pkcs12_path.exists():
            raise ValueError("PKCS#12 certificate file was not found")
        # Ensure the input is parseable and not encrypted before attempting signature.
        _ = _load_pdf(input_path)
        passphrase = (pkcs12_password or "").encode("utf-8") if pkcs12_password else None
        with input_path.open("rb") as source, output_path.open("wb") as sink:
            writer = IncrementalPdfFileWriter(source)
            try:
                signer = signers.SimpleSigner.load_pkcs12(str(pkcs12_path), passphrase=passphrase)
            except TypeError:
                signer = signers.SimpleSigner.load_pkcs12(str(pkcs12_path), passphrase)
            signature_meta = signers.PdfSignatureMetadata(field_name="Signature1")
            pdf_signer = signers.PdfSigner(signature_meta, signer=signer)
            pdf_signer.sign_pdf(writer, output=sink)
        return output_path

    signature_text = text.strip()
    image_path = signature_image_path if signature_image_path and signature_image_path.exists() else None
    if not signature_text and image_path is None:
        raise ValueError("Provide signature text or a signature image for visual signing")

    with fitz.open(str(input_path)) as document:
        _assert_fitz_unencrypted(document)
        target_pages = _resolve_page_selection(
            pages,
            document.page_count,
            tolerant=tolerant_ranges,
        )
        for index in range(document.page_count):
            page_number = index + 1
            if target_pages is not None and page_number not in target_pages:
                continue
            page = document.load_page(index)
            stamp_text = f"Signed: {signature_text}" if signature_text else ""
            text_width = fitz.get_text_length(stamp_text, fontsize=font_size) if stamp_text else 0.0
            base_width = max(160.0, min(text_width + 16.0, 420.0))
            base_height = max(40.0, min(font_size + 22.0, 120.0))
            if image_path is not None:
                base_width = max(base_width, 220.0)
                base_height = max(base_height, 80.0 if stamp_text else 56.0)
            page_rect = page.rect
            anchor_mode = anchor.strip().lower() or "custom"
            supported_anchors = {
                "custom",
                "bottom-right",
                "bottom-left",
                "top-right",
                "top-left",
            }
            if anchor_mode not in supported_anchors:
                raise ValueError(
                    f"Unsupported anchor: {anchor!r}. Use one of: {', '.join(sorted(supported_anchors))}"
                )
            rect = _resolve_signature_rect(page_rect, anchor_mode, base_width, base_height, x, y)
            page.draw_rect(rect, color=(0.22, 0.35, 0.22), width=1)
            if image_path is not None:
                image_bottom = rect.y1 - (font_size + 10 if stamp_text else 4)
                image_rect = fitz.Rect(rect.x0 + 4, rect.y0 + 4, rect.x1 - 4, image_bottom)
                if image_rect.width > 2 and image_rect.height > 2:
                    page.insert_image(image_rect, filename=str(image_path), keep_proportion=True)
            if stamp_text:
                text_x = min(rect.x0 + 8, max(rect.x0, rect.x1 - 8))
                text_y = min(rect.y0 + font_size + 6, max(rect.y0 + 1, rect.y1 - 2))
                page.insert_text(
                    (text_x, text_y),
                    stamp_text,
                    fontsize=font_size,
                    color=(0.1, 0.2, 0.1),
                )
        document.save(
            str(output_path),
            deflate=True,
            garbage=4,
            clean=True,
            incremental=False,
        )
    return output_path


def _parse_edit_operations(raw_operations: object) -> list[dict]:
    """Normalize edit operations into a list of dicts."""
    if raw_operations is None:
        return []
    if isinstance(raw_operations, str):
        text = raw_operations.strip()
        if not text:
            return []
        try:
            parsed = json.loads(text)
        except json.JSONDecodeError as error:
            raise ValueError("Edit operations JSON is invalid") from error
        raw_operations = parsed
    if not isinstance(raw_operations, list):
        raise ValueError("Edit operations must be a list")
    normalized: list[dict] = []
    for item in raw_operations:
        if not isinstance(item, dict):
            raise ValueError("Each edit operation must be an object")
        normalized.append(item)
    return normalized


def _parse_edit_int(value: object, op: str, field: str, default: int) -> int:
    """Parse an integer edit-operation field with a clear error message."""
    if value is None:
        return default
    try:
        return int(value)
    except (TypeError, ValueError) as error:
        raise ValueError(
            f"Invalid numeric value for {field} in {op}: {value!r}"
        ) from error


def _parse_edit_float(value: object, op: str, field: str, default: float) -> float:
    """Parse a float edit-operation field with a clear error message."""
    if value is None:
        return default
    try:
        return float(value)
    except (TypeError, ValueError) as error:
        raise ValueError(
            f"Invalid numeric value for {field} in {op}: {value!r}"
        ) from error


def edit_pdf(
    input_path: Path,
    output_path: Path,
    operations: object,
) -> Path:
    """
    Apply structured edit operations to a PDF in sequence.

    Page references are interpreted against the document state after previous
    operations in the same request.
    """
    normalized_ops = _parse_edit_operations(operations)
    if not normalized_ops:
        raise ValueError("At least one edit operation is required")
    with fitz.open(str(input_path)) as document:
        _assert_fitz_unencrypted(document)
        for operation in normalized_ops:
            op = str(operation.get("op", "")).strip().lower()
            page_number = _parse_edit_int(operation.get("page"), op, "page", 1)
            if op == "delete_pages":
                pages_value = str(operation.get("pages", "")).strip()
                if not pages_value:
                    raise ValueError("delete_pages requires pages")
                try:
                    pages = sorted(
                        _resolve_page_selection(pages_value, document.page_count) or []
                    )
                except ValueError as error:
                    raise ValueError(
                        f"Invalid pages value for delete_pages: {pages_value!r}"
                    ) from error
                for value in reversed(pages):
                    document.delete_page(value - 1)
                continue
            if op == "insert_blank_page":
                index = max(0, min(page_number - 1, document.page_count))
                width = _parse_edit_float(operation.get("width"), op, "width", 612)
                height = _parse_edit_float(operation.get("height"), op, "height", 792)
                document.new_page(pno=index, width=width, height=height)
                continue
            if page_number < 1 or page_number > document.page_count:
                raise ValueError("Edit operation page is out of range")
            page = document.load_page(page_number - 1)
            if op == "add_text":
                text = str(operation.get("text", "")).strip()
                if not text:
                    raise ValueError("add_text requires text")
                x = _parse_edit_float(operation.get("x"), op, "x", 72)
                y = _parse_edit_float(operation.get("y"), op, "y", 72)
                font_size = _parse_edit_float(
                    operation.get("font_size"), op, "font_size", 14
                )
                page.insert_text((x, y), text, fontsize=font_size, color=(0, 0, 0))
                continue
            if op == "draw_rect":
                rect = fitz.Rect(
                    _parse_edit_float(operation.get("x"), op, "x", 72),
                    _parse_edit_float(operation.get("y"), op, "y", 72),
                    _parse_edit_float(operation.get("x2"), op, "x2", 172),
                    _parse_edit_float(operation.get("y2"), op, "y2", 132),
                )
                page.draw_rect(
                    rect,
                    color=(0, 0, 0),
                    width=_parse_edit_float(operation.get("width"), op, "width", 1),
                )
                continue
            if op == "draw_line":
                page.draw_line(
                    (
                        _parse_edit_float(operation.get("x"), op, "x", 72),
                        _parse_edit_float(operation.get("y"), op, "y", 72),
                    ),
                    (
                        _parse_edit_float(operation.get("x2"), op, "x2", 172),
                        _parse_edit_float(operation.get("y2"), op, "y2", 72),
                    ),
                    color=(0, 0, 0),
                    width=_parse_edit_float(operation.get("width"), op, "width", 1),
                )
                continue
            if op == "whiteout":
                rect = fitz.Rect(
                    _parse_edit_float(operation.get("x"), op, "x", 72),
                    _parse_edit_float(operation.get("y"), op, "y", 72),
                    _parse_edit_float(operation.get("x2"), op, "x2", 172),
                    _parse_edit_float(operation.get("y2"), op, "y2", 132),
                )
                page.draw_rect(rect, color=(1, 1, 1), fill=(1, 1, 1), width=0)
                continue
            raise ValueError(f"Unsupported edit op: {op}")
        if document.page_count == 0:
            raise ValueError("Edit operation removed all pages")
        document.save(str(output_path), deflate=True)
    return output_path


def organize_pdf(
    input_path: Path,
    output_path: Path,
    order: str | None = None,
    delete: str | None = None,
    rotate: str | None = None,
    tolerant_ranges: bool = False,
) -> Path:
    """Apply remove, reorder, and rotate in one deterministic operation."""
    reader = _load_pdf(input_path)
    total_pages = len(reader.pages)
    delete_pages_set = _resolve_page_selection(
        delete,
        total_pages,
        tolerant=tolerant_ranges,
    ) or set()
    if order and order.strip():
        ordered_pages: list[int] = []
        seen: set[int] = set()
        for page_number in _parse_page_list(
            order,
            total_pages,
            tolerant=tolerant_ranges,
        ):
            if page_number in delete_pages_set or page_number in seen:
                continue
            ordered_pages.append(page_number)
            seen.add(page_number)
    else:
        ordered_pages = [value for value in range(1, total_pages + 1) if value not in delete_pages_set]
    if not ordered_pages:
        raise ValueError("Organize result has no pages")

    rotations: dict[int, int] = {}
    if rotate and rotate.strip():
        for part in rotate.split(","):
            cleaned = part.strip()
            if not cleaned:
                continue
            if ":" not in cleaned:
                raise ValueError("Rotate format must be page:angle")
            page_text, angle_text = cleaned.split(":", 1)
            try:
                page_value = int(page_text.strip())
                angle_value = int(angle_text.strip())
            except ValueError as error:
                raise ValueError(
                    f"Invalid rotate value {cleaned!r}: page and angle must be integers"
                ) from error
            if page_value < 1 or page_value > total_pages:
                raise ValueError("Rotate page is out of range")
            if angle_value not in (90, 180, 270):
                raise ValueError("Rotate angle must be 90, 180, or 270")
            rotations[page_value] = angle_value

    writer = PdfWriter()
    for page_number in ordered_pages:
        page = reader.pages[page_number - 1]
        angle = rotations.get(page_number)
        if angle:
            _rotate_page(page, angle)
        writer.add_page(page)
    _copy_metadata(writer, reader)
    with output_path.open("wb") as handle:
        writer.write(handle)
    return output_path


def ocr_pdf(
    input_path: Path,
    output_path: Path,
    lang: str | None = None,
    profile: str | None = None,
) -> Path:
    """Convert a PDF into a searchable OCR PDF."""
    language = (lang or DEFAULT_OCR_LANG).strip() or DEFAULT_OCR_LANG
    normalized_profile = _normalize_ocr_profile(profile)
    if os.getenv("ZENPDF_OCR_USE_OCRMYPDF", "1") == "1":
        ocrmypdf = shutil.which("ocrmypdf")
        if ocrmypdf:
            try:
                profile_settings = _ocr_profile_settings(normalized_profile)
                image_dpi = profile_settings["dpi"]
                ocrmypdf_flags: list[str] = []
                if normalized_profile == "fast":
                    ocrmypdf_flags.extend(["--optimize", "0", "--fast-web-view", "1"])
                elif normalized_profile == "accurate":
                    ocrmypdf_flags.extend(["--optimize", "2", "--deskew", "--clean-final"])
                result = subprocess.run(
                    [
                        ocrmypdf,
                        "--skip-text",
                        "--output-type",
                        "pdf",
                        "--language",
                        language,
                        "--image-dpi",
                        str(image_dpi),
                        *ocrmypdf_flags,
                        str(input_path),
                        str(output_path),
                    ],
                    capture_output=True,
                    text=True,
                    check=False,
                    timeout=240,
                )
            except subprocess.TimeoutExpired as error:
                raise RuntimeError("OCR PDF conversion timed out") from error
            if result.returncode == 0 and output_path.exists():
                return output_path
    if pytesseract is None:
        raise RuntimeError("OCR PDF conversion requires pytesseract")
    if (
        not shutil.which("tesseract")
        and getattr(pytesseract, "__name__", "") == "pytesseract"
    ):
        raise RuntimeError("OCR PDF conversion requires tesseract")
    page_outputs: list[Path] = []
    run_id = uuid.uuid4().hex[:8]
    with fitz.open(str(input_path)) as document:
        _assert_fitz_unencrypted(document)
        for index in range(document.page_count):
            page = document.load_page(index)
            image = _render_page_image(page, profile=normalized_profile)
            try:
                pdf_bytes = pytesseract.image_to_pdf_or_hocr(
                    image,
                    extension="pdf",
                    lang=language,
                    timeout=90,
                )
            except TypeError:
                pdf_bytes = pytesseract.image_to_pdf_or_hocr(
                    image,
                    extension="pdf",
                    lang=language,
                )
            page_path = (
                output_path.parent
                / f"{output_path.stem}_ocr_{run_id}_page_{index + 1}.pdf"
            )
            page_path.write_bytes(pdf_bytes)
            page_outputs.append(page_path)
    try:
        merge_pdfs(page_outputs, output_path)
    finally:
        for page_path in page_outputs:
            page_path.unlink(missing_ok=True)
    return output_path


PDF_A_TIMEOUT_SEC = 120
PDF_A_VERSION_TIMEOUT_SEC = 10
PDF_A_MIN_VERSION = (10, 3, 1)


def _parse_version_tuple(raw: str) -> Tuple[int, int, int]:
    """Parse a version string into a comparable tuple."""
    token = raw.strip().split()[0] if raw.strip() else ""
    parts = [part for part in token.split(".") if part]
    numbers: List[int] = []
    for part in parts:
        if not part.isdigit():
            break
        numbers.append(int(part))
    if not numbers:
        raise ValueError("Unable to parse version")
    while len(numbers) < 3:
        numbers.append(0)
    return (numbers[0], numbers[1], numbers[2])


def _verify_pdfa_conformance(output_path: Path) -> dict[str, Any]:
    """Verify PDF/A output and return conformance status details."""
    verapdf = shutil.which("verapdf")
    if verapdf:
        result = _run_command(
            [verapdf, "--format", "text", str(output_path)],
            timeout=60,
            error_prefix="veraPDF validation",
        )
        stdout = (result.stdout or "").lower()
        if result.returncode != 0:
            raise RuntimeError("PDF/A conformance validation failed")
        if "non-compliant" in stdout or "failed" in stdout:
            raise RuntimeError("Output is not PDF/A compliant")
        return {
            "status": "success",
            "standard": "PDF/A-2b",
            "compliant": True,
            "validator": "veraPDF",
            "details": "veraPDF reported compliant output",
        }

    if not _is_valid_pdf(output_path):
        raise RuntimeError("PDF/A output failed structural validation")
    return {
        "status": "success",
        "standard": "PDF/A-2b",
        "compliant": True,
        "validator": "qpdf-structural",
        "details": "veraPDF unavailable; structural PDF validation passed",
    }


def pdf_to_pdfa(
    input_path: Path,
    output_path: Path,
    include_report: bool = False,
) -> Path | tuple[Path, dict[str, Any]]:
    """
    Convert a PDF into PDF/A-2b using Ghostscript.
    
    Parameters:
        input_path (Path): Path to the source PDF file.
        output_path (Path): Destination path for the PDF/A file.
    
    Returns:
        Path: The output_path of the generated PDF/A file.
    
    Raises:
        RuntimeError: If Ghostscript is missing, times out, or fails the conversion.
        ValueError: If the input PDF is encrypted.
    """
    reader = PdfReader(str(input_path))
    if reader.is_encrypted:
        raise ValueError("Encrypted PDFs are not supported for PDF/A conversion")

    ghostscript = shutil.which("gs")
    if not ghostscript:
        raise RuntimeError("Ghostscript is required for PDF/A conversion")

    version_result = subprocess.run(
        [ghostscript, "--version"],
        capture_output=True,
        text=True,
        check=False,
        timeout=PDF_A_VERSION_TIMEOUT_SEC,
    )
    if version_result.returncode != 0:
        raise RuntimeError("Ghostscript version check failed")
    version_output = (version_result.stdout or version_result.stderr or "").strip()
    try:
        version = _parse_version_tuple(version_output)
    except ValueError as error:
        raise RuntimeError("Ghostscript >= 10.03.1 is required for PDF/A conversion") from error
    if version < PDF_A_MIN_VERSION:
        raise RuntimeError("Ghostscript >= 10.03.1 is required for PDF/A conversion")

    output_path.parent.mkdir(parents=True, exist_ok=True)

    command = [
        ghostscript,
        "-dSAFER",
        "-dPDFA=2",
        "-dBATCH",
        "-dNOPAUSE",
        "-dNOOUTERSAVE",
        "-sDEVICE=pdfwrite",
        "-dPDFACompatibilityPolicy=1",
        "-sProcessColorModel=DeviceRGB",
        "-sColorConversionStrategy=RGB",
        "-dUseCIEColor",
        f"-sOutputFile={output_path}",
        str(input_path),
    ]

    try:
        result = subprocess.run(
            command,
            capture_output=True,
            text=True,
            check=False,
            timeout=PDF_A_TIMEOUT_SEC,
        )
    except subprocess.TimeoutExpired as error:
        raise RuntimeError("PDF/A conversion timed out") from error

    if result.returncode != 0:
        raise RuntimeError(result.stderr or result.stdout or "PDF/A conversion failed")

    if not output_path.exists():
        raise RuntimeError("PDF/A conversion produced no output")

    conformance = _verify_pdfa_conformance(output_path)
    if include_report:
        return output_path, conformance
    return output_path


def _extract_text_lines_by_page(input_path: Path) -> List[List[str]]:
    """Return non-empty extracted text lines per page."""
    reader = _load_pdf(input_path)
    lines_by_page: List[List[str]] = []
    for page in reader.pages:
        text = page.extract_text() or ""
        lines = [line.strip() for line in text.splitlines() if line.strip()]
        lines_by_page.append(lines)
    return lines_by_page


def _write_docx_from_pages(lines_by_page: List[List[str]], output_path: Path) -> Path:
    """Write per-page text lines into a DOCX file."""
    document = Document()
    for index, lines in enumerate(lines_by_page, start=1):
        if index > 1:
            document.add_page_break()
        if lines:
            for line in lines:
                document.add_paragraph(line)
        else:
            document.add_paragraph("")
    document.save(str(output_path))
    return output_path


def _extract_tables_by_page(input_path: Path) -> List[List[List[List[str]]]]:
    """Extract tables as page -> table -> row -> cell via pdfplumber."""
    if pdfplumber is None:
        return []
    tables_by_page: List[List[List[List[str]]]] = []
    with pdfplumber.open(str(input_path)) as pdf:
        for page in pdf.pages:
            page_tables: List[List[List[str]]] = []
            for table in page.extract_tables() or []:
                table_rows: List[List[str]] = []
                for raw_row in table or []:
                    if raw_row is None:
                        continue
                    row = [str(cell).strip() if cell is not None else "" for cell in raw_row]
                    if any(value for value in row):
                        table_rows.append(row)
                if table_rows:
                    page_tables.append(table_rows)
            tables_by_page.append(page_tables)
    return tables_by_page


def _write_docx_layout(input_path: Path, output_path: Path) -> Path:
    """Write DOCX using layout-aware text blocks, extracted tables, and page images."""
    document = Document()
    tables_by_page = _extract_tables_by_page(input_path)
    with fitz.open(str(input_path)) as pdf:
        for page_index in range(pdf.page_count):
            if page_index > 0:
                document.add_page_break()
            page = pdf.load_page(page_index)
            emitted = False

            page_tables: List[List[List[str]]] = []
            if page_index < len(tables_by_page):
                page_tables = tables_by_page[page_index]
            for table_rows in page_tables:
                if not table_rows:
                    continue
                col_count = max(len(row) for row in table_rows)
                doc_table = document.add_table(rows=len(table_rows), cols=col_count)
                for row_index, row_values in enumerate(table_rows):
                    for col_index in range(col_count):
                        value = row_values[col_index] if col_index < len(row_values) else ""
                        doc_table.cell(row_index, col_index).text = value
                document.add_paragraph("")
                emitted = True

            page_blocks = page.get_text("dict").get("blocks", [])
            for block in page_blocks:
                block_type = int(block.get("type", -1))
                if block_type == 0:
                    lines: list[str] = []
                    for line in block.get("lines", []):
                        spans = line.get("spans", [])
                        line_text = "".join(str(span.get("text", "")) for span in spans).strip()
                        if line_text:
                            lines.append(line_text)
                    text = "\n".join(lines).strip()
                    if text:
                        document.add_paragraph(text)
                        emitted = True
                elif block_type == 1:
                    image_bytes = block.get("image")
                    if isinstance(image_bytes, (bytes, bytearray)):
                        try:
                            document.add_picture(BytesIO(bytes(image_bytes)))
                            emitted = True
                        except Exception:
                            continue

            if not emitted:
                document.add_paragraph("")
    document.save(str(output_path))
    return output_path


def _ocr_lines_by_page(input_path: Path, lang: str, profile: str | None) -> List[List[str]]:
    """Extract OCR text lines per page."""
    document = fitz.open(str(input_path))
    lines_by_page: List[List[str]] = []
    try:
        for index in range(document.page_count):
            page = document.load_page(index)
            image = _render_page_image(page, profile=profile)
            page_text = _ocr_image(image, lang)
            lines_by_page.append(
                [line.strip() for line in page_text.splitlines() if line.strip()]
            )
    finally:
        document.close()
    return lines_by_page


def pdf_to_docx(
    input_path: Path,
    output_path: Path,
    mode: str = "auto",
    ocr_profile: str | None = None,
) -> Path:
    """Convert PDF to DOCX via layout/text extraction with OCR fallback."""
    normalized_mode = (mode or "auto").strip().lower()
    normalized_profile = _normalize_ocr_profile(ocr_profile)
    if normalized_mode not in {"auto", "layout", "text", "ocr"}:
        raise ValueError("Mode must be auto, layout, text, or ocr")

    if normalized_mode == "ocr":
        lines_by_page = _ocr_lines_by_page(input_path, DEFAULT_OCR_LANG, normalized_profile)
        return _write_docx_from_pages(lines_by_page, output_path)
    if normalized_mode == "text":
        lines_by_page = _extract_text_lines_by_page(input_path)
        return _write_docx_from_pages(lines_by_page, output_path)
    if normalized_mode == "layout":
        return _write_docx_layout(input_path, output_path)
    if _is_text_light_pdf(input_path):
        lines_by_page = _ocr_lines_by_page(input_path, DEFAULT_OCR_LANG, normalized_profile)
        return _write_docx_from_pages(lines_by_page, output_path)
    return _write_docx_layout(input_path, output_path)


def _render_page_image(
    page: fitz.Page,
    dpi: int | None = None,
    profile: str | None = None,
) -> Image.Image:
    """Render a PDF page to a PIL image for OCR."""
    settings = _ocr_profile_settings(profile)
    resolved_dpi = int(dpi or settings["dpi"])
    median_filter_size = int(settings["median_filter"])
    threshold = int(settings["threshold"])
    scale = resolved_dpi / 72
    matrix = fitz.Matrix(scale, scale)
    pix = page.get_pixmap(matrix=matrix, alpha=False)
    image = Image.open(BytesIO(pix.tobytes("png")))
    image = image.convert("L")
    image = ImageOps.autocontrast(image)
    if median_filter_size > 1:
        image = image.filter(ImageFilter.MedianFilter(size=median_filter_size))
    image = image.point(
        lambda px: 255 if px >= threshold else 0, mode="1"
    ).convert("RGB")
    return image


def _ocr_image(image: Image.Image, lang: str) -> str:
    """Run OCR on a PIL image using Tesseract."""
    if pytesseract is None:
        raise RuntimeError("pytesseract is required for OCR conversions")
    if not shutil.which("tesseract"):
        raise RuntimeError("Tesseract is required for OCR conversions")
    safe_lang = (lang or DEFAULT_OCR_LANG).strip() or DEFAULT_OCR_LANG
    return pytesseract.image_to_string(image, lang=safe_lang)


def _extract_table_rows_by_page(input_path: Path) -> List[List[List[str]]]:
    """Backwards-compatible page->rows table extraction view."""
    flattened: List[List[List[str]]] = []
    for page_tables in _extract_tables_by_page(input_path):
        page_rows: List[List[str]] = []
        for table_rows in page_tables:
            page_rows.extend(table_rows)
        flattened.append(page_rows)
    return flattened


def _split_line_into_cells(line: str) -> list[str]:
    """Split a text line into probable tabular cells."""
    stripped = line.strip()
    if not stripped:
        return []
    if "\t" in stripped:
        cells = [cell.strip() for cell in stripped.split("\t")]
    elif "|" in stripped:
        cells = [cell.strip() for cell in stripped.split("|")]
    elif re.search(r"\s{2,}", stripped):
        cells = [cell.strip() for cell in re.split(r"\s{2,}", stripped)]
    elif stripped.count(",") >= 2:
        cells = [cell.strip() for cell in stripped.split(",")]
    else:
        cells = [stripped]
    return [cell for cell in cells if cell]


def _autosize_sheet(sheet, max_columns: int) -> None:  # type: ignore[no-untyped-def]
    """Apply simple width sizing for the first max_columns columns."""
    for col_index in range(1, max_columns + 1):
        max_length = 0
        for row in sheet.iter_rows(min_col=col_index, max_col=col_index):
            value = row[0].value
            if value is None:
                continue
            max_length = max(max_length, len(str(value)))
        sheet.column_dimensions[get_column_letter(col_index)].width = min(64, max(10, max_length + 2))


def _write_xlsx_from_text(lines_by_page: List[List[str]], output_path: Path) -> Path:
    """Write text lines into a structured XLSX workbook."""
    workbook = Workbook()
    default_sheet = workbook.active
    workbook.remove(default_sheet)
    for index, lines in enumerate(lines_by_page, start=1):
        sheet = workbook.create_sheet(title=f"Page_{index}")
        rows = [_split_line_into_cells(line) for line in lines if line.strip()]
        rows = [row for row in rows if row]
        if not rows:
            sheet.cell(row=1, column=1, value=f"Page {index}")
            _autosize_sheet(sheet, 1)
            continue
        max_columns = max(len(row) for row in rows)
        for row_index, row_values in enumerate(rows, start=1):
            for col_index, value in enumerate(row_values, start=1):
                sheet.cell(row=row_index, column=col_index, value=value)
        _autosize_sheet(sheet, max_columns)
    if not workbook.sheetnames:
        workbook.create_sheet("Page_1")
    workbook.save(str(output_path))
    return output_path


def _write_xlsx_from_tables(
    table_rows_by_page: List[List[List[List[str]]]], output_path: Path
) -> Path:
    """Write extracted tables into structured XLSX sheets."""
    workbook = Workbook()
    default_sheet = workbook.active
    workbook.remove(default_sheet)
    created = False
    for page_index, page_tables in enumerate(table_rows_by_page, start=1):
        if not page_tables:
            continue
        for table_index, table_rows in enumerate(page_tables, start=1):
            if not table_rows:
                continue
            sheet = workbook.create_sheet(title=f"P{page_index}_T{table_index}")
            max_columns = max(len(row) for row in table_rows)
            for row_index, row_values in enumerate(table_rows, start=1):
                for col_index in range(1, max_columns + 1):
                    value = row_values[col_index - 1] if col_index <= len(row_values) else ""
                    sheet.cell(row=row_index, column=col_index, value=value)
            _autosize_sheet(sheet, max_columns)
            created = True
    if not created:
        fallback = workbook.create_sheet("Page_1")
        fallback.cell(row=1, column=1, value="No tables detected")
        _autosize_sheet(fallback, 1)
    workbook.save(str(output_path))
    return output_path


def pdf_to_xlsx(
    input_path: Path,
    output_path: Path,
    mode: str = "auto",
    ocr_profile: str | None = None,
) -> Path:
    """Convert PDF to XLSX via table extraction with OCR/text structured fallback."""
    normalized_mode = (mode or "auto").strip().lower()
    normalized_profile = _normalize_ocr_profile(ocr_profile)
    if normalized_mode not in {"auto", "table", "text", "ocr"}:
        raise ValueError("Mode must be auto, table, text, or ocr")

    if normalized_mode == "table":
        tables = _extract_tables_by_page(input_path)
        if tables and any(page_tables for page_tables in tables):
            return _write_xlsx_from_tables(tables, output_path)
        return _write_xlsx_from_text(_extract_text_lines_by_page(input_path), output_path)

    if normalized_mode == "text":
        return _write_xlsx_from_text(_extract_text_lines_by_page(input_path), output_path)

    if normalized_mode == "ocr":
        ocr_lines = _ocr_lines_by_page(input_path, DEFAULT_OCR_LANG, normalized_profile)
        return _write_xlsx_from_text(ocr_lines, output_path)

    tables = _extract_tables_by_page(input_path)
    if tables and any(page_tables for page_tables in tables):
        return _write_xlsx_from_tables(tables, output_path)
    if _is_text_light_pdf(input_path):
        ocr_lines = _ocr_lines_by_page(input_path, DEFAULT_OCR_LANG, normalized_profile)
        return _write_xlsx_from_text(ocr_lines, output_path)
    return _write_xlsx_from_text(_extract_text_lines_by_page(input_path), output_path)


def _resolve_public_ip(hostname: str) -> str:
    """Resolve a hostname to a public IP address."""
    try:
        infos = socket.getaddrinfo(hostname, None)
    except socket.gaierror as error:
        raise ValueError("Unable to resolve host") from error

    public_ips: List[str] = []
    for info in infos:
        address = str(info[4][0])
        if _is_public_ip(address):
            public_ips.append(address)

    if public_ips:
        ipv4_candidates: List[str] = []
        for address in public_ips:
            try:
                ip = ipaddress.ip_address(address)
            except ValueError:
                continue
            if isinstance(ip, ipaddress.IPv4Address):
                ipv4_candidates.append(address)
        if ipv4_candidates:
            return ipv4_candidates[0]
        return public_ips[0]
    raise ValueError("URL host is not allowed")


def _is_public_ip(address: str) -> bool:
    """Return True if the address is a public IP."""
    try:
        ip = ipaddress.ip_address(address)
    except ValueError:
        return False
    return not (
        ip.is_private
        or ip.is_loopback
        or ip.is_link_local
        or ip.is_multicast
        or ip.is_reserved
        or ip.is_unspecified
    )
