"""Tests for worker conversion utilities."""

import os
from pathlib import Path
from tempfile import TemporaryDirectory
from io import BytesIO

import shutil
import subprocess
from unittest.mock import patch

import pytest
import requests
import fitz
from docx import Document
from fpdf import FPDF
from openpyxl import load_workbook
from PIL import Image
from pypdf import PdfReader, PdfWriter

try:
    from pptx import Presentation as PptxPresentation
except ImportError:  # pragma: no cover - optional dependency for PPTX output
    PptxPresentation = None

from zenpdf_worker.tools import (
    MAX_WEB_BYTES,
    OCR_TOOL_TIMEOUT_SECONDS,
    compare_pdfs,
    compress_pdf,
    crop_pdf,
    edit_pdf,
    excel_to_pdf,
    html_to_pdf,
    image_to_pdf,
    merge_pdfs,
    ocr_pdf,
    office_to_pdf,
    organize_pdf,
    page_numbers_pdf,
    pdf_to_pdfa,
    pdf_to_docx,
    pdf_to_powerpoint,
    pdf_to_xlsx,
    pdf_to_jpg,
    powerpoint_to_pdf,
    protect_pdf,
    repair_pdf,
    redact_pdf,
    rotate_pdf,
    scan_to_pdf,
    sign_pdf,
    split_pdf,
    unlock_pdf,
    watermark_pdf,
    web_to_pdf,
    word_to_pdf,
    zip_outputs,
)


def _make_pdf(path: Path, pages: int) -> None:
    """
    Create a blank PDF at the given path containing the specified number of pages.
    
    Parameters:
        path (Path): Destination file path for the generated PDF.
        pages (int): Number of blank pages to include (must be greater than or equal to 0).
    """
    writer = PdfWriter()
    for _ in range(pages):
        writer.add_blank_page(width=300, height=300)
    with path.open("wb") as handle:
        writer.write(handle)


def _make_text_pdf(path: Path, text: str) -> None:
    """Create a single-page PDF with text content."""
    pdf = FPDF(orientation="P", unit="mm", format="A4")
    pdf.add_page()
    pdf.set_font("Helvetica", size=14)
    pdf.text(10, 20, text)
    pdf.output(str(path))


def _make_image_pdf(path: Path) -> None:
    """Create a single-page PDF containing a raster image."""
    image = Image.new("RGB", (300, 300), color=(120, 140, 180))
    with TemporaryDirectory() as temp:
        temp_path = Path(temp)
        image_path = temp_path / "image.png"
        image.save(image_path)
        document = fitz.open()
        page = document.new_page(width=300, height=300)
        page.insert_image(fitz.Rect(0, 0, 300, 300), filename=str(image_path))
        document.save(str(path))
        document.close()


def _make_pdf_with_metadata(path: Path, title: str) -> None:
    """Create a PDF with a Title metadata value."""
    writer = PdfWriter()
    writer.add_blank_page(width=300, height=300)
    writer.add_metadata({"/Title": title})
    with path.open("wb") as handle:
        writer.write(handle)


def test_merge_pdfs() -> None:
    """Merge multiple PDFs into one output."""
    with TemporaryDirectory() as temp:
        temp_path = Path(temp)
        first = temp_path / "first.pdf"
        second = temp_path / "second.pdf"
        _make_pdf(first, 1)
        _make_pdf(second, 2)

        output = merge_pdfs([first, second], temp_path / "merged.pdf")
        reader = PdfReader(str(output))
        assert len(reader.pages) == 3


def test_split_pdf_ranges() -> None:
    """Split a PDF using page ranges."""
    with TemporaryDirectory() as temp:
        temp_path = Path(temp)
        source = temp_path / "source.pdf"
        _make_pdf(source, 3)

        outputs = split_pdf(source, temp_path, "1-2,3")
        assert len(outputs) == 2
        reader = PdfReader(str(outputs[0]))
        assert len(reader.pages) == 2


def test_rotate_pdf() -> None:
    """Rotate pages in a PDF."""
    with TemporaryDirectory() as temp:
        temp_path = Path(temp)
        source = temp_path / "source.pdf"
        _make_pdf(source, 1)
        output = rotate_pdf(source, temp_path / "rotated.pdf", 90, None)
        assert output.exists()


def test_rotate_pdf_rejects_invalid_pages() -> None:
    """Reject malformed page range selections for rotate."""
    with TemporaryDirectory() as temp:
        temp_path = Path(temp)
        source = temp_path / "source.pdf"
        _make_pdf(source, 2)
        with pytest.raises(ValueError):
            rotate_pdf(source, temp_path / "rotated.pdf", 90, "1,,2")


def test_watermark_and_page_numbers() -> None:
    """Apply watermark and page numbers to a PDF."""
    with TemporaryDirectory() as temp:
        temp_path = Path(temp)
        source = temp_path / "source.pdf"
        _make_pdf(source, 2)

        watermarked = watermark_pdf(source, temp_path / "watermarked.pdf", "CONFIDENTIAL", None)
        numbered = page_numbers_pdf(source, temp_path / "numbered.pdf", 3, None)

        watermarked_reader = PdfReader(str(watermarked))
        numbered_reader = PdfReader(str(numbered))
        assert "CONFIDENTIAL" in (watermarked_reader.pages[0].extract_text() or "")
        assert "3" in (numbered_reader.pages[0].extract_text() or "")


def test_page_numbers_selection_mode() -> None:
    """Selection indexing should increment only on selected pages."""
    with TemporaryDirectory() as temp:
        temp_path = Path(temp)
        source = temp_path / "source.pdf"
        _make_pdf(source, 3)
        numbered = page_numbers_pdf(
            source,
            temp_path / "numbered_selection.pdf",
            10,
            "2-3",
            "selectionIndex",
        )
        reader = PdfReader(str(numbered))
        page1 = reader.pages[0].extract_text() or ""
        page2 = reader.pages[1].extract_text() or ""
        page3 = reader.pages[2].extract_text() or ""
        assert "10" not in page1
        assert "10" in page2
        assert "11" in page3


def test_watermark_rejects_invalid_pages() -> None:
    """Reject invalid page selections when watermarking."""
    with TemporaryDirectory() as temp:
        temp_path = Path(temp)
        source = temp_path / "source.pdf"
        _make_pdf(source, 2)

        with pytest.raises(ValueError):
            watermark_pdf(source, temp_path / "watermarked.pdf", "NOTE", "nope")


def test_crop_pdf() -> None:
    """Crop a PDF using point margins."""
    with TemporaryDirectory() as temp:
        temp_path = Path(temp)
        source = temp_path / "source.pdf"
        _make_pdf(source, 1)

        cropped = crop_pdf(source, temp_path / "cropped.pdf", "10,10,10,10", None)
        reader = PdfReader(str(cropped))
        page = reader.pages[0]
        assert float(page.cropbox.width) == pytest.approx(280)
        assert float(page.cropbox.height) == pytest.approx(280)


def test_crop_pdf_rejects_invalid_margins() -> None:
    """Reject invalid margin specifications when cropping."""
    with TemporaryDirectory() as temp:
        temp_path = Path(temp)
        source = temp_path / "source.pdf"
        _make_pdf(source, 1)

        for index, margins in enumerate(["bad", "10,10", "-5"], start=1):
            with pytest.raises(ValueError):
                crop_pdf(
                    source,
                    temp_path / f"invalid_{index}.pdf",
                    margins,
                    None,
                )


def test_image_to_pdf_and_pdf_to_jpg() -> None:
    """Convert image to PDF and PDF to JPG."""
    with TemporaryDirectory() as temp:
        temp_path = Path(temp)
        image_path = temp_path / "sample.png"
        image = Image.new("RGB", (200, 200), color=(120, 140, 180))
        image.save(image_path)

        pdf_path = image_to_pdf([image_path], temp_path / "image.pdf")
        assert pdf_path.exists()

        images = pdf_to_jpg(pdf_path, temp_path, dpi=72)
        assert len(images) == 1
        zipped = zip_outputs(images, temp_path / "pages.zip")
        assert zipped.exists()


def test_unlock_and_protect_pdf() -> None:
    """Protect a PDF with a password and unlock it."""
    with TemporaryDirectory() as temp:
        temp_path = Path(temp)
        source = temp_path / "source.pdf"
        _make_pdf(source, 1)

        protected = protect_pdf(source, temp_path / "protected.pdf", "secret")
        protected_reader = PdfReader(str(protected))
        assert protected_reader.is_encrypted

        unlocked = unlock_pdf(protected, temp_path / "unlocked.pdf", "secret")
        unlocked_reader = PdfReader(str(unlocked))
        assert not unlocked_reader.is_encrypted


def test_repair_pdf() -> None:
    """Rewrite a PDF into a repaired copy."""
    with TemporaryDirectory() as temp:
        temp_path = Path(temp)
        source = temp_path / "source.pdf"
        _make_pdf(source, 2)

        repaired = repair_pdf(source, temp_path / "repaired.pdf")
        reader = PdfReader(str(repaired))
        assert len(reader.pages) == 2


def test_watermark_preserves_metadata() -> None:
    """Keep PDF metadata after watermarking, numbering, and rotating."""
    with TemporaryDirectory() as temp:
        temp_path = Path(temp)
        source = temp_path / "source.pdf"
        _make_pdf_with_metadata(source, "ZenPDF")

        watermarked = watermark_pdf(source, temp_path / "watermarked.pdf", "NOTE", None)
        numbered = page_numbers_pdf(source, temp_path / "numbered.pdf", 1, None)
        rotated = rotate_pdf(source, temp_path / "rotated.pdf", 90, None)

        watermarked_reader = PdfReader(str(watermarked))
        numbered_reader = PdfReader(str(numbered))
        rotated_reader = PdfReader(str(rotated))
        watermarked_metadata = watermarked_reader.metadata or {}
        numbered_metadata = numbered_reader.metadata or {}
        rotated_metadata = rotated_reader.metadata or {}
        assert watermarked_metadata.get("/Title") == "ZenPDF"
        assert numbered_metadata.get("/Title") == "ZenPDF"
        assert rotated_metadata.get("/Title") == "ZenPDF"


def test_redact_pdf() -> None:
    """Redact matching text in a PDF."""
    with TemporaryDirectory() as temp:
        temp_path = Path(temp)
        source = temp_path / "source.pdf"
        _make_text_pdf(source, "CONFIDENTIAL")

        redacted = redact_pdf(source, temp_path / "redacted.pdf", "CONFIDENTIAL", None)
        reader = PdfReader(str(redacted))
        assert "CONFIDENTIAL" not in (reader.pages[0].extract_text() or "")


def test_redact_pdf_whole_word_mode() -> None:
    """Whole-word redaction should not remove longer containing words."""
    with TemporaryDirectory() as temp:
        temp_path = Path(temp)
        source = temp_path / "source.pdf"
        _make_text_pdf(source, "secret secretive")

        redacted = redact_pdf(
            source,
            temp_path / "redacted.pdf",
            "secret",
            None,
            whole_word=True,
        )
        reader = PdfReader(str(redacted))
        page_text = reader.pages[0].extract_text() or ""
        assert "secretive" in page_text


def test_compare_pdfs() -> None:
    """Generate a comparison report for two PDFs."""
    with TemporaryDirectory() as temp:
        temp_path = Path(temp)
        first = temp_path / "first.pdf"
        second = temp_path / "second.pdf"
        _make_text_pdf(first, "Alpha")
        _make_text_pdf(second, "Beta")

        report = compare_pdfs(first, second, temp_path / "report.txt")
        report_text = report.read_text(encoding="utf-8")
        assert "text differs" in report_text
        assert "Visual comparison" in report_text


def test_html_to_pdf() -> None:
    """Render HTML content into a PDF."""
    with TemporaryDirectory() as temp:
        temp_path = Path(temp)
        output = html_to_pdf("<h1>Hello</h1><p>ZenPDF</p>", temp_path / "web.pdf")
        assert output.exists()
        assert output.stat().st_size > 0


def test_web_to_pdf_fetches_html() -> None:
    """Fetch HTML over HTTP and render to PDF."""
    with TemporaryDirectory() as temp:
        temp_path = Path(temp)
        response = _DummyResponse(b"<p>Example</p>")
        session = _DummySession(response)
        with patch("zenpdf_worker.tools.requests.Session", return_value=session), patch(
            "zenpdf_worker.tools._resolve_public_ip", return_value="93.184.216.34"
        ):
            output = web_to_pdf(
                "https://example.com",
                temp_path / "site.pdf",
                render_mode="text",
            )
        assert output.exists()


def test_web_to_pdf_browser_mode() -> None:
    """Browser mode should render through the Playwright path."""
    with TemporaryDirectory() as temp:
        temp_path = Path(temp)
        response = _DummyResponse(b"<p>Example</p>")
        session = _DummySession(response)
        output = temp_path / "site_browser.pdf"
        launch_options: dict[str, object] = {}
        context_options: dict[str, object] = {}

        class _FakePage:
            def __init__(self) -> None:
                self._route_handler = None

            def route(self, _pattern: str, handler) -> None:
                self._route_handler = handler

            def set_content(self, html: str, **_kwargs) -> None:
                assert html == "<p>Example</p>"

            def wait_for_load_state(self, *_args, **_kwargs) -> None:
                return None

            def emulate_media(self, **_kwargs) -> None:
                return None

            def pdf(self, path: str, **_kwargs) -> None:
                Path(path).write_bytes(b"%PDF-1.7\n")

        class _FakeContext:
            def __init__(self) -> None:
                self.page = _FakePage()

            def new_page(self) -> _FakePage:
                return self.page

            def close(self) -> None:
                return None

        class _FakeBrowser:
            def __init__(self) -> None:
                self.context = _FakeContext()

            def new_context(self, **_kwargs) -> _FakeContext:
                context_options.update(_kwargs)
                return self.context

            def close(self) -> None:
                return None

        class _FakePlaywright:
            def __init__(self) -> None:
                class _Chromium:
                    @staticmethod
                    def launch(**_kwargs):
                        launch_options.update(_kwargs)
                        return _FakeBrowser()

                self.chromium = _Chromium()

            def __enter__(self):
                return self

            def __exit__(self, exc_type, exc, tb):
                return False

        with patch("zenpdf_worker.tools.requests.Session", return_value=session), patch(
            "zenpdf_worker.tools._resolve_public_ip", return_value="93.184.216.34"
        ) as resolver:
            with patch("zenpdf_worker.tools.sync_playwright", return_value=_FakePlaywright()):
                rendered = web_to_pdf(
                    "https://example.com",
                    output,
                    render_mode="browser",
                )
        assert rendered.exists()
        assert resolver.call_count == 1
        assert "--host-resolver-rules=MAP * 0.0.0.0" in launch_options["args"]
        assert context_options["java_script_enabled"] is False
        assert context_options["service_workers"] == "block"


def test_web_to_pdf_browser_mode_blocks_private_subresource() -> None:
    """Browser mode should reject blocked private network requests."""
    with TemporaryDirectory() as temp:
        temp_path = Path(temp)
        response = _DummyResponse(b"<p>Example</p>")
        session = _DummySession(response)

        class _DummyRoute:
            def __init__(self) -> None:
                self.aborted = False

            def abort(self) -> None:
                self.aborted = True

            def continue_(self) -> None:
                return None

        class _DummyRequest:
            def __init__(self, url: str) -> None:
                self.url = url

        class _FakePage:
            def __init__(self) -> None:
                self._route_handler = None

            def route(self, _pattern: str, handler) -> None:
                self._route_handler = handler

            def set_content(self, _html: str, **_kwargs) -> None:
                assert self._route_handler is not None
                route = _DummyRoute()
                request = _DummyRequest("http://127.0.0.1/internal")
                self._route_handler(route, request)

            def wait_for_load_state(self, *_args, **_kwargs) -> None:
                return None

            def emulate_media(self, **_kwargs) -> None:
                return None

            def pdf(self, path: str, **_kwargs) -> None:
                Path(path).write_bytes(b"%PDF-1.7\n")

        class _FakeContext:
            def __init__(self) -> None:
                self.page = _FakePage()

            def new_page(self) -> _FakePage:
                return self.page

            def close(self) -> None:
                return None

        class _FakeBrowser:
            def __init__(self) -> None:
                self.context = _FakeContext()

            def new_context(self, **_kwargs) -> _FakeContext:
                return self.context

            def close(self) -> None:
                return None

        class _FakePlaywright:
            def __init__(self) -> None:
                class _Chromium:
                    @staticmethod
                    def launch(**_kwargs):
                        return _FakeBrowser()

                self.chromium = _Chromium()

            def __enter__(self):
                return self

            def __exit__(self, exc_type, exc, tb):
                return False

        def _resolve_ip(host: str) -> str:
            if host == "127.0.0.1":
                raise ValueError("URL host is not allowed")
            return "93.184.216.34"

        with patch("zenpdf_worker.tools.requests.Session", return_value=session), patch(
            "zenpdf_worker.tools._resolve_public_ip", side_effect=_resolve_ip
        ):
            with patch("zenpdf_worker.tools.sync_playwright", return_value=_FakePlaywright()):
                with pytest.raises(ValueError, match="Blocked private"):
                    web_to_pdf(
                        "https://example.com",
                        temp_path / "site_browser.pdf",
                        render_mode="browser",
                    )


def test_split_pdf_tolerant_ranges() -> None:
    """Tolerant range mode should process valid tokens and skip invalid tokens."""
    with TemporaryDirectory() as temp:
        temp_path = Path(temp)
        source = temp_path / "source.pdf"
        _make_pdf(source, 3)
        outputs = split_pdf(source, temp_path, "1-2,bad,3", tolerant_ranges=True)
        assert len(outputs) == 2


def test_sign_pdf_visual_image_mode() -> None:
    """Visual signing should accept an uploaded signature image."""
    with TemporaryDirectory() as temp:
        temp_path = Path(temp)
        source = temp_path / "source.pdf"
        image_path = temp_path / "signature.png"
        _make_pdf(source, 1)
        Image.new("RGB", (140, 48), color=(40, 40, 40)).save(image_path)
        signed = sign_pdf(
            source,
            temp_path / "signed_image.pdf",
            text="",
            signature_image_path=image_path,
        )
        assert signed.exists()


def test_pdfa_include_report(monkeypatch: pytest.MonkeyPatch) -> None:
    """PDF/A conversion should return conformance metadata when requested."""
    with TemporaryDirectory() as temp:
        temp_path = Path(temp)
        source = temp_path / "source.pdf"
        output = temp_path / "output.pdf"
        _make_pdf(source, 1)

        def fake_which(name: str):
            if name == "gs":
                return "/usr/bin/gs"
            if name == "verapdf":
                return None
            return None

        monkeypatch.setattr("zenpdf_worker.tools.shutil.which", fake_which)

        def fake_run(command, **_kwargs):
            if "--version" in command:
                return subprocess.CompletedProcess(command, 0, stdout="10.03.1", stderr="")
            writer = PdfWriter()
            writer.add_blank_page(width=300, height=300)
            with output.open("wb") as handle:
                writer.write(handle)
            return subprocess.CompletedProcess(command, 0, stdout="", stderr="")

        monkeypatch.setattr("zenpdf_worker.tools.subprocess.run", fake_run)

        converted, result = pdf_to_pdfa(source, output, include_report=True)
        assert converted.exists()
        assert result.get("compliant") is True


def test_pdf_to_powerpoint_editable_mode() -> None:
    """Editable mode should produce PPTX output or surface dependency errors."""
    with TemporaryDirectory() as temp:
        temp_path = Path(temp)
        source = temp_path / "source.pdf"
        _make_text_pdf(source, "Editable")
        output = temp_path / "editable.pptx"
        try:
            result = pdf_to_powerpoint(source, output, mode="editable")
        except RuntimeError as error:
            assert "python-pptx is required" in str(error) or "utilities" in str(error)
            return
        assert result.exists()


def test_web_to_pdf_blocks_private_host() -> None:
    """Reject private hosts for web-to-pdf."""
    with TemporaryDirectory() as temp:
        temp_path = Path(temp)
        with patch(
            "zenpdf_worker.tools._resolve_public_ip",
            side_effect=ValueError("URL host is not allowed"),
        ):
            with pytest.raises(ValueError):
                web_to_pdf(
                    "http://127.0.0.1",
                    temp_path / "blocked.pdf",
                    render_mode="text",
                )


def test_web_to_pdf_blocks_redirects() -> None:
    """Reject redirect responses for web-to-pdf."""
    with TemporaryDirectory() as temp:
        temp_path = Path(temp)
        response = _DummyResponse(b"", status_code=302)
        session = _DummySession(response)
        with patch("zenpdf_worker.tools.requests.Session", return_value=session), patch(
            "zenpdf_worker.tools._resolve_public_ip", return_value="93.184.216.34"
        ):
            with pytest.raises(ValueError):
                web_to_pdf(
                    "https://example.com",
                    temp_path / "redirect.pdf",
                    render_mode="text",
                )


def test_web_to_pdf_never_falls_back_to_hostname() -> None:
    """An IP-bound TLS failure must not trigger a second DNS-based fetch."""
    with TemporaryDirectory() as temp:
        temp_path = Path(temp)
        response = _DummyResponse(b"<p>Example</p>")

        class _FailingSession(_DummySession):
            def get(self, *_args, **_kwargs):
                raise requests.exceptions.SSLError("handshake failed")

        class _SessionFactory:
            def __init__(self):
                self.calls = 0

            def __call__(self):
                self.calls += 1
                return _FailingSession(response)

        factory = _SessionFactory()

        with patch("zenpdf_worker.tools.requests.Session", side_effect=factory), patch(
            "zenpdf_worker.tools._resolve_public_ip", return_value="93.184.216.34"
        ) as resolver, patch.dict(
            os.environ,
            {"ZENPDF_WEB_ALLOW_HOSTNAME_FALLBACK": "1", "ZENPDF_DEV_MODE": "1"},
        ):
            with pytest.raises(requests.exceptions.SSLError):
                web_to_pdf(
                    "https://example.com",
                    temp_path / "site.pdf",
                    render_mode="text",
                )

        assert resolver.call_count == 1
        assert factory.calls == 1


def test_web_to_pdf_binds_fetch_to_single_dns_resolution() -> None:
    """The actual request must use the one validated IP, not resolve the hostname again."""
    with TemporaryDirectory() as temp:
        temp_path = Path(temp)
        response = _DummyResponse(b"<p>Example</p>")

        class _RecordingSession(_DummySession):
            def __init__(self) -> None:
                super().__init__(response)
                self.urls: list[str] = []

            def get(self, url: str, *_args, **_kwargs):
                self.urls.append(url)
                return super().get(url, *_args, **_kwargs)

        session = _RecordingSession()
        with patch("zenpdf_worker.tools.requests.Session", return_value=session), patch(
            "zenpdf_worker.tools._resolve_public_ip",
            side_effect=["93.184.216.34", "127.0.0.1"],
        ) as resolver:
            output = web_to_pdf(
                "https://example.com/document",
                temp_path / "site.pdf",
                render_mode="text",
            )

        assert output.exists()
        assert resolver.call_count == 1
        assert session.urls == ["https://93.184.216.34/document"]


def test_web_to_pdf_limits_body_size() -> None:
    """Enforce the max response size for web-to-pdf."""
    with TemporaryDirectory() as temp:
        temp_path = Path(temp)
        response = _DummyResponse(b"a" * (MAX_WEB_BYTES + 1))
        session = _DummySession(response)
        with patch("zenpdf_worker.tools.requests.Session", return_value=session), patch(
            "zenpdf_worker.tools._resolve_public_ip", return_value="93.184.216.34"
        ):
            with pytest.raises(ValueError):
                web_to_pdf(
                    "https://example.com",
                    temp_path / "large.pdf",
                    render_mode="text",
                )


def test_pdf_to_docx_and_xlsx() -> None:
    """Convert PDF text to DOCX and XLSX files in text mode."""
    with TemporaryDirectory() as temp:
        temp_path = Path(temp)
        source = temp_path / "source.pdf"
        _make_text_pdf(source, "Hello ZenPDF")

        docx_path = pdf_to_docx(source, temp_path / "output.docx", mode="text")
        xlsx_path = pdf_to_xlsx(source, temp_path / "output.xlsx", mode="text")

        assert docx_path.exists()
        assert xlsx_path.exists()
        document = Document(str(docx_path))
        doc_text = "\n".join(paragraph.text for paragraph in document.paragraphs)
        assert "Hello ZenPDF" in doc_text
        workbook = load_workbook(xlsx_path)
        values = [cell.value for cell in workbook.active["A"] if cell.value]
        assert "Hello ZenPDF" in values


def test_pdf_to_docx_and_xlsx_ocr_modes() -> None:
    """Convert a PDF to DOCX/XLSX using OCR mode."""
    with TemporaryDirectory() as temp:
        temp_path = Path(temp)
        source = temp_path / "source.pdf"
        _make_pdf(source, 2)

        with patch(
            "zenpdf_worker.tools._ocr_image",
            side_effect=["First page", "Second page"],
        ):
            docx_path = pdf_to_docx(source, temp_path / "ocr.docx", mode="ocr")

        with patch(
            "zenpdf_worker.tools._ocr_image",
            side_effect=["First sheet", "Second sheet"],
        ):
            xlsx_path = pdf_to_xlsx(source, temp_path / "ocr.xlsx", mode="ocr")

        document = Document(str(docx_path))
        doc_text = "\n".join(paragraph.text for paragraph in document.paragraphs)
        assert "First page" in doc_text
        assert "Second page" in doc_text

        workbook = load_workbook(xlsx_path)
        values: list[str] = []
        for sheet in workbook.worksheets:
            values.extend(str(cell.value) for cell in sheet["A"] if cell.value)
        assert "First sheet" in values
        assert "Second sheet" in values


def test_pdf_to_docx_auto_uses_ocr_for_text_light(monkeypatch: pytest.MonkeyPatch) -> None:
    """Auto mode should choose OCR for text-light PDFs."""
    with TemporaryDirectory() as temp:
        temp_path = Path(temp)
        source = temp_path / "source.pdf"
        _make_pdf(source, 1)
        monkeypatch.setattr("zenpdf_worker.tools._is_text_light_pdf", lambda _path: True)
        with patch("zenpdf_worker.tools._ocr_image", return_value="OCR line"):
            docx_path = pdf_to_docx(source, temp_path / "auto.docx", mode="auto")
        document = Document(str(docx_path))
        values = "\n".join(paragraph.text for paragraph in document.paragraphs)
        assert "OCR line" in values


def test_compress_pdf_detects_image_heavy(monkeypatch: pytest.MonkeyPatch) -> None:
    """Classify an image-heavy PDF and return compression metadata."""
    with TemporaryDirectory() as temp:
        temp_path = Path(temp)
        source = temp_path / "image.pdf"
        output = temp_path / "output.pdf"
        _make_image_pdf(source)

        monkeypatch.setattr(shutil, "which", lambda _tool: None)
        monkeypatch.setenv("ZENPDF_COMPRESS_SAVINGS_THRESHOLD_PCT", "0")
        monkeypatch.setenv("ZENPDF_COMPRESS_MIN_SAVINGS_BYTES", "0")
        monkeypatch.setenv("ZENPDF_COMPRESS_AUTO_IMAGE_HEAVY", "1")
        monkeypatch.setenv("ZENPDF_COMPRESS_TIMEOUT_SECONDS", "5")

        _, result = compress_pdf(source, output)
        assert result["image_metrics"]["image_heavy"] is True


def test_compress_pdf_detects_text_heavy(monkeypatch: pytest.MonkeyPatch) -> None:
    """Classify a text-heavy PDF and return compression metadata."""
    with TemporaryDirectory() as temp:
        temp_path = Path(temp)
        source = temp_path / "text.pdf"
        output = temp_path / "output.pdf"
        _make_text_pdf(source, "Hello ZenPDF")

        monkeypatch.setattr(shutil, "which", lambda _tool: None)
        monkeypatch.setenv("ZENPDF_COMPRESS_SAVINGS_THRESHOLD_PCT", "0")
        monkeypatch.setenv("ZENPDF_COMPRESS_MIN_SAVINGS_BYTES", "0")
        monkeypatch.setenv("ZENPDF_COMPRESS_AUTO_IMAGE_HEAVY", "1")
        monkeypatch.setenv("ZENPDF_COMPRESS_TIMEOUT_SECONDS", "5")

        _, result = compress_pdf(source, output)
        assert result["image_metrics"]["image_heavy"] is False


def test_compress_pdf_rejects_encrypted_pdf(monkeypatch: pytest.MonkeyPatch) -> None:
    """Reject encrypted PDFs before compression."""
    with TemporaryDirectory() as temp:
        temp_path = Path(temp)
        source = temp_path / "encrypted.pdf"
        output = temp_path / "output.pdf"
        writer = PdfWriter()
        writer.add_blank_page(width=300, height=300)
        writer.encrypt("secret")
        with source.open("wb") as handle:
            writer.write(handle)

        monkeypatch.setattr(shutil, "which", lambda _tool: None)
        with pytest.raises(ValueError, match="PDF is encrypted"):
            compress_pdf(source, output)


def test_pdfa_requires_ghostscript(monkeypatch: pytest.MonkeyPatch) -> None:
    """Fail when Ghostscript is not available for PDF/A conversion."""
    with TemporaryDirectory() as temp:
        temp_path = Path(temp)
        source = temp_path / "source.pdf"
        _make_pdf(source, 1)

        monkeypatch.setattr("zenpdf_worker.tools.shutil.which", lambda _: None)
        with pytest.raises(RuntimeError):
            pdf_to_pdfa(source, temp_path / "output.pdf")


def test_pdfa_conversion_runs_ghostscript(monkeypatch: pytest.MonkeyPatch) -> None:
    """Run PDF/A conversion through Ghostscript."""
    with TemporaryDirectory() as temp:
        temp_path = Path(temp)
        source = temp_path / "source.pdf"
        output = temp_path / "output.pdf"
        _make_pdf(source, 1)

        def fake_which(name: str):
            if name == "gs":
                return "/usr/bin/gs"
            return None

        monkeypatch.setattr("zenpdf_worker.tools.shutil.which", fake_which)

        def fake_run(command, **_kwargs):
            if "--version" in command:
                return subprocess.CompletedProcess(command, 0, stdout="10.03.1", stderr="")
            writer = PdfWriter()
            writer.add_blank_page(width=300, height=300)
            with output.open("wb") as handle:
                writer.write(handle)
            return subprocess.CompletedProcess(command, 0, stdout="", stderr="")

        monkeypatch.setattr("zenpdf_worker.tools.subprocess.run", fake_run)

        result = pdf_to_pdfa(source, output)
        assert result.exists()


def test_pdfa_rejects_encrypted_pdf() -> None:
    """Reject encrypted PDFs before invoking Ghostscript."""
    with TemporaryDirectory() as temp:
        temp_path = Path(temp)
        source = temp_path / "encrypted.pdf"
        writer = PdfWriter()
        writer.add_blank_page(width=300, height=300)
        writer.encrypt("secret")
        with source.open("wb") as handle:
            writer.write(handle)

        with pytest.raises(ValueError):
            pdf_to_pdfa(source, temp_path / "output.pdf")


def test_office_to_pdf_missing_soffice(monkeypatch: pytest.MonkeyPatch) -> None:
    """Fail when LibreOffice is not available."""
    with TemporaryDirectory() as temp:
        temp_path = Path(temp)
        doc_path = temp_path / "sample.docx"
        document = Document()
        document.add_paragraph("Hello")
        document.save(doc_path)

        monkeypatch.setattr("zenpdf_worker.tools.shutil.which", lambda _: None)
        with pytest.raises(RuntimeError):
            office_to_pdf(doc_path, temp_path)


def test_word_powerpoint_excel_to_pdf_extension_validation() -> None:
    """Reject unsupported source extensions for split Office conversion tools."""
    with TemporaryDirectory() as temp:
        temp_path = Path(temp)
        source = temp_path / "sample.txt"
        source.write_text("not office", encoding="utf-8")
        with pytest.raises(ValueError):
            word_to_pdf(source, temp_path)
        with pytest.raises(ValueError):
            powerpoint_to_pdf(source, temp_path)
        with pytest.raises(ValueError):
            excel_to_pdf(source, temp_path)


def test_organize_pdf_order_delete_rotate() -> None:
    """Apply organize operation in deterministic delete->order->rotate behavior."""
    with TemporaryDirectory() as temp:
        temp_path = Path(temp)
        source = temp_path / "source.pdf"
        _make_text_pdf(source, "one")
        # build 3-page input
        page_two = temp_path / "two.pdf"
        page_three = temp_path / "three.pdf"
        _make_text_pdf(page_two, "two")
        _make_text_pdf(page_three, "three")
        merge_pdfs([source, page_two, page_three], source)

        output = organize_pdf(
            source,
            temp_path / "organized.pdf",
            order="3,1,2",
            delete="2",
            rotate="3:90",
        )
        reader = PdfReader(str(output))
        assert len(reader.pages) == 2


def test_edit_pdf_add_text_and_delete_pages() -> None:
    """Apply edit operations including text insertion and page deletion."""
    with TemporaryDirectory() as temp:
        temp_path = Path(temp)
        source = temp_path / "source.pdf"
        _make_pdf(source, 2)
        operations = [
            {"op": "add_text", "page": 1, "x": 72, "y": 72, "text": "Approved"},
            {"op": "delete_pages", "pages": "2"},
        ]
        edited = edit_pdf(source, temp_path / "edited.pdf", operations)
        reader = PdfReader(str(edited))
        assert len(reader.pages) == 1
        assert "Approved" in (reader.pages[0].extract_text() or "")


def test_sign_pdf_adds_signature_stamp() -> None:
    """Apply visible signature text to selected pages."""
    with TemporaryDirectory() as temp:
        temp_path = Path(temp)
        source = temp_path / "source.pdf"
        _make_pdf(source, 1)
        signed = sign_pdf(source, temp_path / "signed.pdf", "Jane Doe")
        reader = PdfReader(str(signed))
        assert "Signed: Jane Doe" in (reader.pages[0].extract_text() or "")


def test_scan_to_pdf_alias() -> None:
    """Scan-to-PDF should produce a merged PDF from images."""
    with TemporaryDirectory() as temp:
        temp_path = Path(temp)
        image_path = temp_path / "scan.png"
        image = Image.new("RGB", (200, 200), color=(120, 140, 180))
        image.save(image_path)
        output = scan_to_pdf([image_path], temp_path / "scan.pdf")
        assert output.exists()


def test_ocr_pdf_fallback(monkeypatch: pytest.MonkeyPatch) -> None:
    """Fallback OCR path should create a searchable PDF when ocrmypdf is unavailable."""
    with TemporaryDirectory() as temp:
        temp_path = Path(temp)
        source = temp_path / "source.pdf"
        _make_image_pdf(source)

        monkeypatch.setenv("ZENPDF_OCR_USE_OCRMYPDF", "0")

        class _FakePytesseract:
            @staticmethod
            def image_to_pdf_or_hocr(
                _image, extension="pdf", lang="eng", timeout=None
            ):
                assert extension == "pdf"
                assert lang == "eng"
                assert timeout == OCR_TOOL_TIMEOUT_SECONDS
                writer = PdfWriter()
                writer.add_blank_page(width=300, height=300)
                buffer = BytesIO()
                writer.write(buffer)
                return buffer.getvalue()

        with patch("zenpdf_worker.tools.pytesseract", _FakePytesseract):
            output = ocr_pdf(source, temp_path / "ocr.pdf", "eng")
        assert output.exists()
        assert output.stat().st_size > 0


def test_pdf_to_powerpoint() -> None:
    """Convert a PDF into PPTX (or surface dependency requirement)."""
    with TemporaryDirectory() as temp:
        temp_path = Path(temp)
        source = temp_path / "source.pdf"
        _make_pdf(source, 1)
        output = temp_path / "slides.pptx"
        try:
            result = pdf_to_powerpoint(source, output)
        except RuntimeError as error:
            assert "python-pptx is required" in str(error)
            return
        assert result.exists()
        assert result.stat().st_size > 0


def test_pdf_to_powerpoint_preserves_page_aspect_ratio() -> None:
    """Use slide dimensions from the source PDF and avoid image distortion."""
    with TemporaryDirectory() as temp:
        temp_path = Path(temp)
        source = temp_path / "wide.pdf"
        writer = PdfWriter()
        writer.add_blank_page(width=640, height=360)
        with source.open("wb") as handle:
            writer.write(handle)

        output = temp_path / "wide_slides.pptx"
        try:
            result = pdf_to_powerpoint(source, output)
        except RuntimeError as error:
            assert "python-pptx is required" in str(error)
            return
        assert result.exists()
        assert result.stat().st_size > 0
        if PptxPresentation is None:
            return

        presentation = PptxPresentation(str(result))
        assert len(presentation.slides) == 1
        expected_ratio = 640 / 360
        slide_ratio = presentation.slide_width / presentation.slide_height
        assert slide_ratio == pytest.approx(expected_ratio, rel=0.01)

        slide = presentation.slides[0]
        assert len(slide.shapes) == 1
        picture = slide.shapes[0]
        picture_ratio = picture.width / picture.height
        assert picture_ratio == pytest.approx(expected_ratio, rel=0.02)


class _DummySocket:
    """Mock socket for response peer IP retrieval."""
    def __init__(self, ip: str) -> None:
        """Initialize the mock socket with an IP."""
        self._ip = ip

    def getpeername(self):
        """Return the mock socket peer tuple."""
        return (self._ip, 443)


class _DummyConnection:
    """Mock connection wrapper."""
    def __init__(self, ip: str) -> None:
        """Initialize the dummy connection with the socket."""
        self.sock = _DummySocket(ip)


class _DummyRaw:
    """Mock raw response wrapper."""
    def __init__(self, ip: str) -> None:
        """Initialize the raw wrapper with a connection."""
        self._connection = _DummyConnection(ip)


class _DummyResponse:
    """Mock response used for web-to-pdf tests."""
    def __init__(self, body: bytes, status_code: int = 200, ip: str = "93.184.216.34") -> None:
        """Initialize the dummy response payload."""
        self._body = body
        self.status_code = status_code
        self.encoding = "utf-8"
        self.raw = _DummyRaw(ip)

    def iter_content(self, _chunk_size: int = 1024, **_kwargs):
        """Yield the body payload as a single chunk."""
        yield self._body

    def raise_for_status(self) -> None:
        """Raise when the status indicates an error."""
        if self.status_code >= 400:
            raise RuntimeError("Bad response")

    def __enter__(self):
        """Enter the response context manager."""
        return self

    def __exit__(self, exc_type, exc, tb):
        """Exit the response context manager."""
        return False


class _DummySession:
    """Mock requests session wrapper."""
    def __init__(self, response: _DummyResponse) -> None:
        """Initialize the dummy session."""
        self._response = response

    def __enter__(self):
        """Enter the session context manager."""
        return self

    def __exit__(self, exc_type, exc, tb):
        """Exit the session context manager."""
        return False

    def close(self):
        """Close the dummy session."""
        return None

    def mount(self, *_args, **_kwargs):
        """Ignore adapter mounting for the mock session."""
        return None

    def get(self, *_args, **_kwargs):
        """Return the prepared response."""
        return self._response
